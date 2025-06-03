import os
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook
import pandas as pd
import json
import xml.etree.ElementTree as ET
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
import numpy as np
from datetime import datetime
import glob
import re
from pathlib import Path

class Enhanced3DPowerPointExtractor:
    def __init__(self, pptx_path=None, user_login="Rateb1223"):
        # Dynamic current time and user
        self.current_time = datetime(2025, 6, 3, 16, 16, 32)  # Current UTC time
        self.user_login = user_login
        
        # Auto-detect PowerPoint file if not provided
        if pptx_path is None:
            pptx_path = self._find_powerpoint_file()
        
        self.pptx_path = pptx_path
        
        # Dynamic output directory with timestamp and user
        base_name = Path(pptx_path).stem if pptx_path else "powerpoint"
        timestamp = self.current_time.strftime("%Y%m%d_%H%M%S")
        self.output_dir = f"enhanced_extraction_{base_name}_{self.user_login}_{timestamp}"
        self.unzipped_dir = f"temp_unzipped_{timestamp}"
        
        # Create structured output directories
        self.dirs = {
            'main': self.output_dir,
            'charts_3d': os.path.join(self.output_dir, "charts_3d"),
            'charts_2d': os.path.join(self.output_dir, "charts_2d"),
            'tables': os.path.join(self.output_dir, "tables"),
            'excel_data': os.path.join(self.output_dir, "excel_data"),
            'xml_data': os.path.join(self.output_dir, "xml_analysis"),
            'slide_text': os.path.join(self.output_dir, "slide_text"),
            'errors': os.path.join(self.output_dir, "errors"),
            'analysis': os.path.join(self.output_dir, "analysis")
        }
        
        for dir_path in self.dirs.values():
            os.makedirs(dir_path, exist_ok=True)
        
        # Initialize counters and storage
        self.chart_counter = 0
        self.chart_3d_counter = 0
        self.table_counter = 0
        self.excel_counter = 0
        self.embedded_data = []
        self.xml_charts = {}
        self.chart_errors = []
        self.table_errors = []
        self.all_charts = []
        self.all_3d_charts = []
        
        print(f"🚀 Enhanced 3D PowerPoint Extractor")
        print(f"👤 User: {self.user_login}")
        print(f"📅 Current Time: {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC")
        print(f"📄 Processing: {os.path.basename(self.pptx_path) if self.pptx_path else 'Auto-detect'}")
        print(f"📂 Output Directory: {self.output_dir}")
    
    def _find_powerpoint_file(self):
        """Auto-detect PowerPoint files in current directory"""
        patterns = ["*.pptx", "*.ppt", "*.potx"]
        files = []
        
        for pattern in patterns:
            files.extend(glob.glob(pattern))
        
        if files:
            # Prefer .pptx files
            pptx_files = [f for f in files if f.endswith('.pptx')]
            if pptx_files:
                selected_file = pptx_files[0]
            else:
                selected_file = files[0]
            
            print(f"🔍 Auto-detected PowerPoint file: {selected_file}")
            return selected_file
        else:
            print("❌ No PowerPoint files found in current directory")
            return None
    
    def extract_all(self):
        """Main extraction method with enhanced 3D chart capabilities"""
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            print(f"❌ PowerPoint file not found: {self.pptx_path}")
            return None
        
        print(f"\n{'='*60}")
        print(f"🔄 Starting Enhanced Extraction Process")
        print(f"{'='*60}")
        
        try:
            # Step 1: Extract and analyze embedded data
            self._extract_embedded_excel_data()
            
            # Step 2: Extract XML chart data for 3D analysis
            self._extract_xml_chart_data()
            
            # Step 3: Load and process PowerPoint presentation
            self._process_presentation()
            
            # Step 4: Generate comprehensive analysis
            self._generate_analysis_reports()
            
            # Step 5: Cleanup temporary files
            self._cleanup()
            
            print(f"\n✅ Enhanced extraction complete!")
            print(f"📊 Total charts: {self.chart_counter}")
            print(f"🎲 3D charts: {self.chart_3d_counter}")
            print(f"📋 Tables: {self.table_counter}")
            print(f"📁 Excel files: {self.excel_counter}")
            print(f"📂 Output: {self.output_dir}")
            
            return self.output_dir
            
        except Exception as e:
            print(f"❌ Error during extraction: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _extract_embedded_excel_data(self):
        """Enhanced embedded Excel data extraction with metadata"""
        print(f"\n📊 Extracting embedded Excel data...")
        
        try:
            # Unzip PowerPoint file
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                zip_ref.extractall(self.unzipped_dir)
            
            # Find embedded Excel files
            embedded_excel_dir = os.path.join(self.unzipped_dir, "ppt", "embeddings")
            
            if not os.path.exists(embedded_excel_dir):
                print("⚠️  No embedded Excel files found")
                return
            
            embedded_excels = sorted([f for f in os.listdir(embedded_excel_dir) 
                                    if f.endswith((".xlsx", ".xls"))])
            
            print(f"📁 Found {len(embedded_excels)} embedded Excel files")
            
            # Process each Excel file
            for file_idx, file in enumerate(embedded_excels):
                try:
                    self.excel_counter += 1
                    file_path = os.path.join(embedded_excel_dir, file)
                    
                    print(f"  📄 Processing Excel {self.excel_counter}: {file}")
                    
                    # Load workbook
                    wb = load_workbook(file_path, data_only=True)
                    
                    excel_data = {
                        'file_name': file,
                        'file_index': self.excel_counter,
                        'extraction_metadata': {
                            'extracted_by': self.user_login,
                            'extraction_time': self.current_time.isoformat(),
                            'file_size': os.path.getsize(file_path)
                        },
                        'sheets': []
                    }
                    
                    # Process each worksheet
                    for sheet_idx, sheet in enumerate(wb.worksheets):
                        print(f"    📋 Sheet: {sheet.title}")
                        
                        # Extract all data
                        rows = []
                        for row in sheet.iter_rows(values_only=True):
                            row_data = [str(cell).strip() if cell is not None else "" for cell in row]
                            rows.append(row_data)
                        
                        # Analyze data structure
                        sheet_analysis = self._analyze_excel_sheet(rows, sheet.title)
                        
                        sheet_data = {
                            'sheet_name': sheet.title,
                            'sheet_index': sheet_idx,
                            'rows': len(rows),
                            'columns': len(rows[0]) if rows else 0,
                            'data': rows,
                            'analysis': sheet_analysis
                        }
                        
                        excel_data['sheets'].append(sheet_data)
                        
                        # Save individual sheet as CSV
                        if rows:
                            csv_path = os.path.join(self.dirs['excel_data'], 
                                                  f"excel_{self.excel_counter}_sheet_{sheet_idx}_{sheet.title}.csv")
                            df = pd.DataFrame(rows[1:], columns=rows[0] if rows else None)
                            df.to_csv(csv_path, index=False, encoding='utf-8')
                    
                    # Save Excel metadata
                    json_path = os.path.join(self.dirs['excel_data'], f"excel_{self.excel_counter}_metadata.json")
                    with open(json_path, 'w', encoding='utf-8') as f:
                        json.dump(excel_data, f, indent=2, ensure_ascii=False)
                    
                    self.embedded_data.append(excel_data)
                    
                except Exception as e:
                    error_info = {
                        'file': file,
                        'error': str(e),
                        'error_type': 'excel_extraction'
                    }
                    self.chart_errors.append(error_info)
                    print(f"    ⚠️  Error processing Excel file: {e}")
        
        except Exception as e:
            print(f"❌ Error extracting embedded Excel data: {e}")
    
    def _analyze_excel_sheet(self, rows, sheet_name):
        """Analyze Excel sheet structure for better chart matching"""
        analysis = {
            'has_headers': False,
            'numeric_columns': [],
            'text_columns': [],
            'potential_categories': [],
            'potential_values': [],
            'data_structure': 'unknown'
        }
        
        if not rows or len(rows) < 2:
            return analysis
        
        # Analyze headers
        headers = rows[0]
        data_rows = rows[1:]
        
        if headers and any(isinstance(h, str) and h.strip() for h in headers):
            analysis['has_headers'] = True
        
        # Analyze column types
        for col_idx, header in enumerate(headers):
            column_values = [row[col_idx] if col_idx < len(row) else '' for row in data_rows]
            
            # Check if column is numeric
            numeric_count = 0
            for val in column_values:
                try:
                    float(str(val).replace(',', ''))
                    numeric_count += 1
                except:
                    pass
            
            if numeric_count > len(column_values) * 0.7:  # 70% numeric
                analysis['numeric_columns'].append(col_idx)
                analysis['potential_values'].extend(column_values)
            else:
                analysis['text_columns'].append(col_idx)
                if col_idx == 0:  # First column often contains categories
                    analysis['potential_categories'] = column_values
        
        # Determine data structure
        if len(analysis['text_columns']) >= 1 and len(analysis['numeric_columns']) >= 1:
            analysis['data_structure'] = 'chart_suitable'
        elif len(analysis['numeric_columns']) > 1:
            analysis['data_structure'] = 'multi_series'
        else:
            analysis['data_structure'] = 'simple_table'
        
        return analysis
    
    def _extract_xml_chart_data(self):
        """Extract XML chart data for advanced 3D analysis"""
        print(f"\n🔍 Extracting XML chart data for 3D analysis...")
        
        try:
            # Find chart XML files
            charts_dir = os.path.join(self.unzipped_dir, "ppt", "charts")
            
            if not os.path.exists(charts_dir):
                print("⚠️  No chart XML files found")
                return
            
            chart_files = [f for f in os.listdir(charts_dir) if f.endswith('.xml')]
            print(f"📊 Found {len(chart_files)} chart XML files")
            
            for chart_file in chart_files:
                try:
                    chart_path = os.path.join(charts_dir, chart_file)
                    
                    with open(chart_path, 'r', encoding='utf-8') as f:
                        xml_content = f.read()
                    
                    # Parse XML for 3D properties
                    chart_analysis = self._analyze_chart_xml(xml_content, chart_file)
                    
                    chart_id = chart_file.replace('.xml', '')
                    self.xml_charts[chart_id] = chart_analysis
                    
                    # Save XML analysis
                    analysis_path = os.path.join(self.dirs['xml_data'], f"{chart_id}_analysis.json")
                    with open(analysis_path, 'w', encoding='utf-8') as f:
                        json.dump(chart_analysis, f, indent=2, ensure_ascii=False)
                    
                    print(f"  📄 Analyzed: {chart_file} (3D: {chart_analysis.get('is_3d', False)})")
                    
                except Exception as e:
                    print(f"    ⚠️  Error analyzing {chart_file}: {e}")
        
        except Exception as e:
            print(f"❌ Error extracting XML chart data: {e}")
    
    def _analyze_chart_xml(self, xml_content, filename):
        """Comprehensive XML chart analysis for 3D detection"""
        analysis = {
            'filename': filename,
            'is_3d': False,
            'chart_type': 'unknown',
            '3d_properties': {},
            'series_count': 0,
            'categories_count': 0,
            'extraction_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            }
        }
        
        try:
            root = ET.fromstring(xml_content)
            
            # Determine namespace
            namespace = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
            if '}' in root.tag:
                ns_uri = root.tag.split('}')[0].strip('{')
                namespace = {'c': ns_uri}
            
            # Find chart type
            chart_type_elem = root.find('.//c:chart', namespace)
            if chart_type_elem is not None:
                # Look for plot area
                plot_area = chart_type_elem.find('.//c:plotArea', namespace)
                if plot_area is not None:
                    # Check different chart types
                    chart_types = [
                        'bar3DChart', 'line3DChart', 'pie3DChart', 'area3DChart',
                        'surface3DChart', 'bubbleChart', 'barChart', 'lineChart'
                    ]
                    
                    for chart_type in chart_types:
                        type_elem = plot_area.find(f'.//c:{chart_type}', namespace)
                        if type_elem is not None:
                            analysis['chart_type'] = chart_type
                            analysis['is_3d'] = '3d' in chart_type.lower() or 'surface' in chart_type.lower()
                            break
            
            # Extract 3D view properties if 3D chart
            if analysis['is_3d']:
                view3d = root.find('.//c:view3D', namespace)
                if view3d is not None:
                    properties = {}
                    
                    # Rotation properties
                    for prop in ['rotX', 'rotY', 'depthPercent', 'rAngAx', 'perspective']:
                        elem = view3d.find(f'./c:{prop}', namespace)
                        if elem is not None and 'val' in elem.attrib:
                            properties[prop] = elem.attrib['val']
                    
                    analysis['3d_properties'] = properties
            
            # Count series and categories
            series_elems = root.findall('.//c:ser', namespace)
            analysis['series_count'] = len(series_elems)
            
            cat_elems = root.findall('.//c:cat', namespace)
            analysis['categories_count'] = len(cat_elems)
            
        except Exception as e:
            analysis['xml_error'] = str(e)
        
        return analysis
    
    def _process_presentation(self):
        """Enhanced PowerPoint presentation processing with 3D chart focus"""
        print(f"\n📖 Processing PowerPoint presentation...")
        
        try:
            prs = Presentation(self.pptx_path)
            excel_idx = 0
            
            print(f"📄 Found {len(prs.slides)} slides to process")
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                print(f"\n🔍 Processing Slide {slide_idx}...")
                
                slide_data = {
                    'slide_number': slide_idx,
                    'slide_title': self._extract_slide_title(slide, slide_idx),
                    'text_content': [],
                    'charts': [],
                    'tables': [],
                    'extraction_metadata': {
                        'extracted_by': self.user_login,
                        'extraction_time': self.current_time.isoformat()
                    }
                }
                
                slide_text_lines = [f"[SLIDE {slide_idx}: {slide_data['slide_title']}]"]
                slide_text_lines.append(f"Extracted by: {self.user_login}")
                slide_text_lines.append(f"Extraction time: {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC")
                slide_text_lines.append("=" * 60)
                
                # Process each shape
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Extract text content
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            text = shape.text_frame.text.strip()
                            slide_data['text_content'].append(text)
                            slide_text_lines.append(f"\n[TEXT {shape_idx + 1}]")
                            slide_text_lines.append(text)
                        
                        # Process tables with enhanced extraction
                        if shape.has_table:
                            table_data = self._process_table_enhanced(shape, slide_idx, shape_idx)
                            if table_data:
                                slide_data['tables'].append(table_data)
                                slide_text_lines.append(f"\n[TABLE {len(slide_data['tables'])}]")
                                slide_text_lines.extend(table_data['text_representation'])
                        
                        # Process charts with enhanced 3D detection
                        if hasattr(shape, 'chart'):
                            chart_data = self._process_chart_enhanced(shape, slide_idx, excel_idx)
                            if chart_data:
                                self.chart_counter += 1
                                slide_data['charts'].append(chart_data)
                                
                                if chart_data.get('is_3d'):
                                    self.chart_3d_counter += 1
                                    self.all_3d_charts.append(chart_data)
                                    slide_text_lines.append(f"\n[3D CHART {self.chart_3d_counter}]")
                                else:
                                    slide_text_lines.append(f"\n[CHART {self.chart_counter}]")
                                
                                slide_text_lines.extend(chart_data['text_representation'])
                                excel_idx += 1
                                
                                self.all_charts.append(chart_data)
                    
                    except Exception as e:
                        error_info = {
                            'slide': slide_idx,
                            'shape': shape_idx,
                            'error': str(e),
                            'error_type': 'shape_processing'
                        }
                        self.chart_errors.append(error_info)
                        print(f"    ⚠️  Error processing shape {shape_idx}: {e}")
                
                # Save slide data
                slide_json_path = os.path.join(self.dirs['slide_text'], f"slide_{slide_idx}_data.json")
                with open(slide_json_path, 'w', encoding='utf-8') as f:
                    json.dump(slide_data, f, indent=2, ensure_ascii=False)
                
                # Save slide text file
                slide_text_path = os.path.join(self.dirs['slide_text'], f"slide_{slide_idx}.txt")
                with open(slide_text_path, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(slide_text_lines))
                
                print(f"  📄 Slide {slide_idx}: {len(slide_data['charts'])} charts, {len(slide_data['tables'])} tables")
        
        except Exception as e:
            print(f"❌ Error processing presentation: {e}")
    
    def _extract_slide_title(self, slide, slide_num):
        """Extract slide title with fallback"""
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except:
            pass
        
        # Look for title-like text
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                if len(text) < 100 and '\n' not in text:
                    return text
        
        return f"Slide {slide_num}"
    
    def _process_table_enhanced(self, table_shape, slide_num, shape_idx):
        """Enhanced table processing with comprehensive data extraction"""
        try:
            self.table_counter += 1
            table_id = f"table_{slide_num}_{self.table_counter}"
            
            print(f"    📊 Table found: {table_id}")
            
            table = table_shape.table
            table_data = {
                'table_id': table_id,
                'slide_number': slide_num,
                'shape_index': shape_idx,
                'rows': len(table.rows),
                'columns': len(table.columns) if len(table.rows) > 0 else 0,
                'headers': [],
                'data': [],
                'text_representation': [],
                'extraction_metadata': {
                    'extracted_by': self.user_login,
                    'extraction_time': self.current_time.isoformat()
                }
            }
            
            # Extract table data
            text_lines = ["TABLE START"]
            
            if len(table.rows) > 0:
                # Headers
                headers = []
                for cell in table.rows[0].cells:
                    header_text = cell.text.strip()
                    headers.append(header_text)
                
                table_data['headers'] = headers
                text_lines.append("\t".join(headers))
                
                # Data rows
                for row in table.rows[1:]:
                    row_data = {}
                    row_text = []
                    
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        row_text.append(cell_text)
                        
                        key = headers[col_idx] if col_idx < len(headers) else f"Column {col_idx + 1}"
                        row_data[key] = cell_text
                    
                    table_data['data'].append(row_data)
                    text_lines.append("\t".join(row_text))
            
            text_lines.append("TABLE END")
            table_data['text_representation'] = text_lines
            
            # Save table as JSON and CSV
            table_json_path = os.path.join(self.dirs['tables'], f"{table_id}.json")
            with open(table_json_path, 'w', encoding='utf-8') as f:
                json.dump(table_data, f, indent=2, ensure_ascii=False)
            
            if table_data['data']:
                table_csv_path = os.path.join(self.dirs['tables'], f"{table_id}.csv")
                df = pd.DataFrame(table_data['data'])
                df.to_csv(table_csv_path, index=False, encoding='utf-8')
            
            return table_data
        
        except Exception as e:
            error_info = {
                'slide': slide_num,
                'shape': shape_idx,
                'error': str(e),
                'error_type': 'table_processing'
            }
            self.table_errors.append(error_info)
            print(f"    ⚠️  Error processing table: {e}")
            return None
    
    def _process_chart_enhanced(self, chart_shape, slide_num, excel_idx):
        """Enhanced chart processing with comprehensive 3D detection and analysis"""
        try:
            chart = chart_shape.chart
            chart_id = f"chart_{slide_num}_{self.chart_counter + 1}"
            
            # Initialize chart data structure
            chart_data = {
                'chart_id': chart_id,
                'slide_number': slide_num,
                'chart_type': str(chart.chart_type),
                'is_3d': False,
                '3d_properties': {},
                'excel_data': None,
                'categories': [],
                'series_data': [],
                'text_representation': [],
                'extraction_metadata': {
                    'extracted_by': self.user_login,
                    'extraction_time': self.current_time.isoformat(),
                    'excel_index': excel_idx
                }
            }
            
            # Enhanced 3D detection
            chart_type_str = str(chart.chart_type).lower()
            chart_data['is_3d'] = any(keyword in chart_type_str for keyword in 
                                    ['3d', 'three', 'dimensional', 'surface', 'cone', 'pyramid'])
            
            # Get XML chart analysis if available
            xml_chart_id = f"chart{self.chart_counter + 1}"
            if xml_chart_id in self.xml_charts:
                xml_analysis = self.xml_charts[xml_chart_id]
                chart_data['is_3d'] = chart_data['is_3d'] or xml_analysis.get('is_3d', False)
                chart_data['3d_properties'].update(xml_analysis.get('3d_properties', {}))
            
            print(f"    {'🎲' if chart_data['is_3d'] else '📊'} Chart: {chart_id} ({'3D' if chart_data['is_3d'] else '2D'})")
            
            # Match with embedded Excel data
            excel_data = None
            if excel_idx < len(self.embedded_data):
                excel_file = self.embedded_data[excel_idx]
                if excel_file['sheets']:
                    excel_data = excel_file['sheets'][0]  # Use first sheet
                    chart_data['excel_data'] = excel_data
                    
                    # Extract enhanced chart data from Excel
                    self._extract_chart_data_from_excel(chart_data, excel_data)
            
            # Create text representation
            text_lines = [f"CHART START - {chart_data['chart_type']}"]
            if chart_data['is_3d']:
                text_lines.append("🎲 3D CHART DETECTED")
                
                # Add 3D properties
                if chart_data['3d_properties']:
                    text_lines.append("3D Properties:")
                    for prop, value in chart_data['3d_properties'].items():
                        text_lines.append(f"  {prop}: {value}")
            
            # Add data information
            if chart_data['categories']:
                text_lines.append(f"Categories: {len(chart_data['categories'])}")
                text_lines.append("Category: Value (Percentage)")
                
                if chart_data['series_data']:
                    for series in chart_data['series_data']:
                        series_name = series.get('name', 'Unnamed Series')
                        values = series.get('values', [])
                        
                        if values and chart_data['categories']:
                            total = sum(values) if sum(values) > 0 else 1
                            
                            for cat, val in zip(chart_data['categories'], values):
                                pct = round((val / total) * 100, 2) if val else 0
                                formatted_val = f"{int(val):,}" if val >= 1000 else f"{val}"
                                text_lines.append(f"{cat}: {formatted_val} ({pct}%)")
            
            text_lines.append("CHART END")
            chart_data['text_representation'] = text_lines
            
            # Create chart visualization
            self._create_chart_visualization(chart_data)
            
            # Save chart data
            chart_json_path = os.path.join(
                self.dirs['charts_3d'] if chart_data['is_3d'] else self.dirs['charts_2d'],
                f"{chart_id}.json"
            )
            with open(chart_json_path, 'w', encoding='utf-8') as f:
                json.dump(chart_data, f, indent=2, ensure_ascii=False)
            
            return chart_data
        
        except Exception as e:
            error_info = {
                'slide': slide_num,
                'chart_type': str(getattr(chart_shape, 'chart', {}).get('chart_type', 'Unknown')),
                'error': str(e),
                'error_type': 'chart_processing'
            }
            self.chart_errors.append(error_info)
            print(f"    ⚠️  Error processing chart: {e}")
            return None
    
    def _extract_chart_data_from_excel(self, chart_data, excel_data):
        """Enhanced Excel data extraction for charts with 3D support"""
        try:
            rows = excel_data.get('data', [])
            analysis = excel_data.get('analysis', {})
            
            if not rows or len(rows) < 2:
                return
            
            headers = rows[0]
            data_rows = rows[1:]
            
            # Extract categories (usually first column)
            if analysis.get('potential_categories'):
                chart_data['categories'] = analysis['potential_categories'][:20]  # Limit to 20
            elif data_rows:
                chart_data['categories'] = [row[0] if row else f"Item {i+1}" 
                                          for i, row in enumerate(data_rows)]
            
            # Extract series data
            numeric_columns = analysis.get('numeric_columns', [])
            
            for col_idx in numeric_columns:
                if col_idx < len(headers):
                    series_name = headers[col_idx] if headers[col_idx] else f"Series {col_idx + 1}"
                    
                    # Extract values
                    values = []
                    for row in data_rows:
                        if col_idx < len(row):
                            try:
                                val = float(str(row[col_idx]).replace(',', ''))
                            except:
                                val = 0
                            values.append(val)
                        else:
                            values.append(0)
                    
                    series_data = {
                        'name': series_name,
                        'values': values,
                        'column_index': col_idx
                    }
                    
                    # For 3D charts, add Z-values if possible
                    if chart_data['is_3d']:
                        series_data['z_values'] = self._generate_3d_values(values, chart_data['chart_type'])
                    
                    chart_data['series_data'].append(series_data)
        
        except Exception as e:
            chart_data['excel_extraction_error'] = str(e)
    
    def _generate_3d_values(self, values, chart_type):
        """Generate appropriate Z-values for 3D charts"""
        if not values:
            return []
        
        chart_type_lower = chart_type.lower()
        
        strategies = {
            'column3d': lambda v: [1.0] * len(v),
            'bar3d': lambda v: [1.0] * len(v),
            'line3d': lambda v: [float(i) * 0.5 for i in range(len(v))],
            'pie3d': lambda v: [0.5] * len(v),
            'area3d': lambda v: [float(val) * 0.1 for val in v],
            'surface3d': lambda v: [float(val) * 0.2 for val in v],
        }
        
        for key, strategy in strategies.items():
            if key in chart_type_lower:
                return strategy(values)
        
        # Default strategy
        return [float(i) + (val * 0.01) for i, val in enumerate(values)]
    
    def _create_chart_visualization(self, chart_data):
        """Create enhanced chart visualizations with 3D support"""
        try:
            chart_id = chart_data['chart_id']
            is_3d = chart_data['is_3d']
            
            if is_3d:
                # Create 3D visualization
                fig = plt.figure(figsize=(12, 9))
                ax = fig.add_subplot(111, projection='3d')
                
                self._create_3d_chart(ax, chart_data)
                
                # Apply 3D properties
                self._apply_3d_properties(ax, chart_data.get('3d_properties', {}))
                
                chart_path = os.path.join(self.dirs['charts_3d'], f"{chart_id}.png")
            else:
                # Create 2D visualization
                fig, ax = plt.subplots(figsize=(10, 6))
                
                self._create_2d_chart(ax, chart_data)
                
                chart_path = os.path.join(self.dirs['charts_2d'], f"{chart_id}.png")
            
            # Add title and metadata
            title = f"{'🎲 3D' if is_3d else '📊'} Chart: {chart_data['chart_type']}"
            ax.set_title(title, fontsize=14, fontweight='bold', pad=20)
            
            # Add extraction metadata
            metadata_text = f"Extracted by {self.user_login} | {self.current_time.strftime('%Y-%m-%d %H:%M')}"
            if is_3d:
                ax.text2D(0.02, 0.02, metadata_text, transform=ax.transAxes, fontsize=8, alpha=0.7)
            else:
                ax.text(0.02, 0.02, metadata_text, transform=ax.transAxes, fontsize=8, alpha=0.7)
            
            plt.tight_layout()
            plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close()
            
            chart_data['visualization_path'] = chart_path
            print(f"      💾 Saved visualization: {os.path.basename(chart_path)}")
        
        except Exception as e:
            print(f"      ⚠️  Error creating visualization: {e}")
    
    def _create_3d_chart(self, ax, chart_data):
        """Create 3D chart visualization"""
        categories = chart_data.get('categories', [])
        series_data = chart_data.get('series_data', [])
        chart_type = chart_data.get('chart_type', '').lower()
        
        if not series_data:
            ax.text2D(0.5, 0.5, "No Data Available", transform=ax.transAxes, 
                     ha='center', va='center', fontsize=14)
            return
        
        # Create appropriate 3D visualization based on chart type
        if 'column' in chart_type or 'bar' in chart_type:
            self._create_3d_bar_chart(ax, categories, series_data)
        elif 'line' in chart_type:
            self._create_3d_line_chart(ax, categories, series_data)
        elif 'pie' in chart_type:
            self._create_3d_pie_chart(ax, categories, series_data)
        elif 'surface' in chart_type or 'area' in chart_type:
            self._create_3d_surface_chart(ax, categories, series_data)
        else:
            # Default to 3D bar
            self._create_3d_bar_chart(ax, categories, series_data)
        
        # Set axis labels
        ax.set_xlabel('Categories')
        ax.set_ylabel('Series')
        ax.set_zlabel('Values')
    
    def _create_3d_bar_chart(self, ax, categories, series_data):
        """Create 3D bar chart"""
        if not categories:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        colors = plt.cm.viridis(np.linspace(0, 1, len(series_data)))
        
        x_pos = np.arange(len(categories))
        y_pos = np.arange(len(series_data))
        
        width = 0.6
        depth = 0.6
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                for j, val in enumerate(values):
                    if j < len(x_pos):
                        x = x_pos[j]
                        y = y_pos[i] if i < len(y_pos) else i
                        z = 0
                        
                        dx = width
                        dy = depth
                        dz = max(float(val), 0.1) if isinstance(val, (int, float)) else 0.1
                        
                        ax.bar3d(x, y, z, dx, dy, dz, color=colors[i], alpha=0.8)
        
        # Set ticks
        ax.set_xticks(x_pos + width/2)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.set_yticks(y_pos + depth/2)
        ax.set_yticklabels([s.get('name', f'Series {i+1}') for i, s in enumerate(series_data)])
    
    def _create_3d_line_chart(self, ax, categories, series_data):
        """Create 3D line chart"""
        colors = plt.cm.tab10(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            z_values = series.get('z_values', [])
            
            if values:
                x = np.arange(len(values))
                y = z_values if z_values and len(z_values) == len(values) else [i] * len(values)
                z = [float(v) if isinstance(v, (int, float)) else 0 for v in values]
                
                ax.plot(x, y, z, color=colors[i], linewidth=3, marker='o', markersize=5,
                       label=series.get('name', f'Series {i+1}'))
        
        ax.legend()
    
    def _create_3d_pie_chart(self, ax, categories, series_data):
        """Create 3D pie chart representation"""
        if not series_data or not series_data[0].get('values'):
            return
        
        values = series_data[0]['values']
        labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
        
        # Create simple 3D representation
        angles = np.linspace(0, 2*np.pi, len(values), endpoint=False)
        colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
        
        for i, (angle, value, color) in enumerate(zip(angles, values, colors)):
            if isinstance(value, (int, float)) and value > 0:
                # Create wedge
                theta = np.linspace(angle, angle + 2*np.pi/len(values), 20)
                r = 3
                height = value / max(values) if max(values) > 0 else 0.1
                
                x = r * np.cos(theta)
                y = r * np.sin(theta)
                z = np.full_like(x, height)
                
                ax.plot(x, y, z, color=color, linewidth=3)
    
    def _create_3d_surface_chart(self, ax, categories, series_data):
        """Create 3D surface chart"""
        if not series_data:
            return
        
        x_size = len(categories) if categories else len(series_data[0].get('values', []))
        y_size = len(series_data)
        
        if x_size == 0 or y_size == 0:
            return
        
        x = np.arange(x_size)
        y = np.arange(y_size)
        X, Y = np.meshgrid(x, y)
        
        Z = np.zeros_like(X, dtype=float)
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            for j, val in enumerate(values):
                if j < Z.shape[1] and isinstance(val, (int, float)):
                    Z[i, j] = float(val)
        
        surf = ax.plot_surface(X, Y, Z, cmap='viridis', alpha=0.8)
        plt.colorbar(surf, ax=ax, shrink=0.5, aspect=5)
    
    def _create_2d_chart(self, ax, chart_data):
        """Create 2D chart visualization"""
        categories = chart_data.get('categories', [])
        series_data = chart_data.get('series_data', [])
        chart_type = chart_data.get('chart_type', '').lower()
        
        if not series_data:
            ax.text(0.5, 0.5, "No Data Available", transform=ax.transAxes,
                   ha='center', va='center', fontsize=14)
            return
        
        # Create appropriate 2D visualization
        if 'bar' in chart_type or 'column' in chart_type:
            self._create_2d_bar_chart(ax, categories, series_data)
        elif 'line' in chart_type:
            self._create_2d_line_chart(ax, categories, series_data)
        elif 'pie' in chart_type:
            self._create_2d_pie_chart(ax, categories, series_data)
        else:
            self._create_2d_bar_chart(ax, categories, series_data)
    
    def _create_2d_bar_chart(self, ax, categories, series_data):
        """Create 2D bar chart"""
        if not categories:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = np.arange(len(categories))
        width = 0.8 / len(series_data) if len(series_data) > 1 else 0.6
        colors = plt.cm.Set3(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                offset = (i - len(series_data)/2 + 0.5) * width
                ax.bar(x + offset, values, width, label=series.get('name', f'Series {i+1}'),
                      color=colors[i], alpha=0.8)
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        
        if len(series_data) > 1:
            ax.legend()
    
    def _create_2d_line_chart(self, ax, categories, series_data):
        """Create 2D line chart"""
        if not categories:
            categories = [f"Point {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = range(len(categories))
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                ax.plot(x, values, marker='o', label=series.get('name', f'Series {i+1}'),
                       color=colors[i], linewidth=2, markersize=4)
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.grid(True, alpha=0.3)
        
        if len(series_data) > 1:
            ax.legend()
    
    def _create_2d_pie_chart(self, ax, categories, series_data):
        """Create 2D pie chart"""
        if series_data and series_data[0].get('values'):
            values = series_data[0]['values']
            labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
            
            # Filter positive values
            filtered_data = [(label, float(value)) for label, value in zip(labels, values)
                           if isinstance(value, (int, float)) and value > 0]
            
            if filtered_data:
                labels, values = zip(*filtered_data)
                colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
                wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%',
                                                 colors=colors, startangle=90)
                
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontweight('bold')
        
        ax.axis('equal')
    
    def _apply_3d_properties(self, ax, properties):
        """Apply 3D view properties to chart"""
        try:
            # Default view
            elev, azim = 20, -35
            
            # Apply rotation if available
            if 'rotX' in properties:
                elev = int(properties['rotX']) / 60000  # Convert from XML units
            if 'rotY' in properties:
                azim = int(properties['rotY']) / 60000
            
            ax.view_init(elev=elev, azim=azim)
            
        except Exception as e:
            # Use default view
            ax.view_init(elev=20, azim=-35)
    
    def _generate_analysis_reports(self):
        """Generate comprehensive analysis reports"""
        print(f"\n📊 Generating analysis reports...")
        
        # Main summary report
        self._generate_main_summary()
        
        # 3D charts analysis
        if self.all_3d_charts:
            self._generate_3d_charts_analysis()
        
        # Excel data analysis
        if self.embedded_data:
            self._generate_excel_analysis()
        
        # Error reports
        if self.chart_errors or self.table_errors:
            self._generate_error_reports()
    
    def _generate_main_summary(self):
        """Generate main summary report"""
        summary_data = {
            'extraction_summary': {
                'source_file': os.path.basename(self.pptx_path),
                'source_path': os.path.abspath(self.pptx_path),
                'extracted_by': self.user_login,
                'extraction_time': self.current_time.isoformat(),
                'extractor_version': 'enhanced_3d_v2.0'
            },
            'statistics': {
                'total_charts': self.chart_counter,
                'charts_3d': self.chart_3d_counter,
                'charts_2d': self.chart_counter - self.chart_3d_counter,
                'tables': self.table_counter,
                'excel_files': self.excel_counter,
                'xml_charts_analyzed': len(self.xml_charts),
                'errors': len(self.chart_errors) + len(self.table_errors)
            },
            'charts_breakdown': {
                'chart_types': {},
                '3d_chart_types': {},
                'charts_with_excel_data': 0
            }
        }
        
        # Analyze chart types
        for chart in self.all_charts:
            chart_type = chart.get('chart_type', 'Unknown')
            summary_data['charts_breakdown']['chart_types'][chart_type] = \
                summary_data['charts_breakdown']['chart_types'].get(chart_type, 0) + 1
            
            if chart.get('is_3d'):
                summary_data['charts_breakdown']['3d_chart_types'][chart_type] = \
                    summary_data['charts_breakdown']['3d_chart_types'].get(chart_type, 0) + 1
            
            if chart.get('excel_data'):
                summary_data['charts_breakdown']['charts_with_excel_data'] += 1
        
        # Save summary
        summary_path = os.path.join(self.dirs['analysis'], 'extraction_summary.json')
        with open(summary_path, 'w', encoding='utf-8') as f:
            json.dump(summary_data, f, indent=2, ensure_ascii=False)
        
        # Create markdown summary
        md_path = os.path.join(self.dirs['analysis'], 'extraction_summary.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(f"# Enhanced PowerPoint Extraction Summary\n\n")
            f.write(f"**Extracted by:** {self.user_login}\n")
            f.write(f"**Extraction Time:** {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC\n")
            f.write(f"**Source File:** {os.path.basename(self.pptx_path)}\n\n")
            
            f.write(f"## 📊 Statistics\n\n")
            f.write(f"- **Total Charts:** {self.chart_counter}\n")
            f.write(f"- **3D Charts:** {self.chart_3d_counter}\n")
            f.write(f"- **2D Charts:** {self.chart_counter - self.chart_3d_counter}\n")
            f.write(f"- **Tables:** {self.table_counter}\n")
            f.write(f"- **Excel Files:** {self.excel_counter}\n")
            f.write(f"- **XML Charts Analyzed:** {len(self.xml_charts)}\n\n")
            
            if self.chart_3d_counter > 0:
                f.write(f"## 🎲 3D Charts Found\n\n")
                for i, chart in enumerate(self.all_3d_charts, 1):
                    f.write(f"{i}. **{chart['chart_id']}** (Slide {chart['slide_number']})\n")
                    f.write(f"   - Type: {chart['chart_type']}\n")
                    if chart.get('3d_properties'):
                        f.write(f"   - 3D Properties: {len(chart['3d_properties'])} found\n")
                    f.write(f"\n")
        
        print(f"  📄 Main summary saved: {md_path}")
    
    def _generate_3d_charts_analysis(self):
        """Generate detailed 3D charts analysis"""
        analysis_data = {
            'total_3d_charts': self.chart_3d_counter,
            'analysis_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            },
            'charts': self.all_3d_charts,
            '3d_properties_summary': {},
            'chart_types_distribution': {}
        }
        
        # Analyze 3D properties
        all_properties = {}
        for chart in self.all_3d_charts:
            chart_type = chart.get('chart_type', 'Unknown')
            analysis_data['chart_types_distribution'][chart_type] = \
                analysis_data['chart_types_distribution'].get(chart_type, 0) + 1
            
            for prop, value in chart.get('3d_properties', {}).items():
                if prop not in all_properties:
                    all_properties[prop] = []
                all_properties[prop].append(value)
        
        analysis_data['3d_properties_summary'] = {
            prop: {'count': len(values), 'unique_values': list(set(values))}
            for prop, values in all_properties.items()
        }
        
        # Save 3D analysis
        analysis_path = os.path.join(self.dirs['analysis'], '3d_charts_analysis.json')
        with open(analysis_path, 'w', encoding='utf-8') as f:
            json.dump(analysis_data, f, indent=2, ensure_ascii=False)
        
        print(f"  🎲 3D charts analysis saved: {analysis_path}")
    
    def _generate_excel_analysis(self):
        """Generate Excel data analysis"""
        excel_analysis = {
            'total_files': self.excel_counter,
            'analysis_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            },
            'files': self.embedded_data,
            'data_structure_summary': {},
            'sheets_summary': {
                'total_sheets': 0,
                'chart_suitable_sheets': 0,
                'multi_series_sheets': 0
            }
        }
        
        # Analyze data structures
        for excel_file in self.embedded_data:
            for sheet in excel_file.get('sheets', []):
                excel_analysis['sheets_summary']['total_sheets'] += 1
                
                structure = sheet.get('analysis', {}).get('data_structure', 'unknown')
                excel_analysis['data_structure_summary'][structure] = \
                    excel_analysis['data_structure_summary'].get(structure, 0) + 1
                
                if structure == 'chart_suitable':
                    excel_analysis['sheets_summary']['chart_suitable_sheets'] += 1
                elif structure == 'multi_series':
                    excel_analysis['sheets_summary']['multi_series_sheets'] += 1
        
        # Save Excel analysis
        analysis_path = os.path.join(self.dirs['analysis'], 'excel_data_analysis.json')
        with open(analysis_path, 'w', encoding='utf-8') as f:
            json.dump(excel_analysis, f, indent=2, ensure_ascii=False)
        
        print(f"  📊 Excel analysis saved: {analysis_path}")
    
    def _generate_error_reports(self):
        """Generate error reports"""
        if self.chart_errors:
            errors_df = pd.DataFrame(self.chart_errors)
            errors_path = os.path.join(self.dirs['errors'], 'chart_errors.xlsx')
            errors_df.to_excel(errors_path, index=False)
            print(f"  ⚠️  Chart errors saved: {errors_path}")
        
        if self.table_errors:
            errors_df = pd.DataFrame(self.table_errors)
            errors_path = os.path.join(self.dirs['errors'], 'table_errors.xlsx')
            errors_df.to_excel(errors_path, index=False)
            print(f"  ⚠️  Table errors saved: {errors_path}")
    
    def _cleanup(self):
        """Clean up temporary files"""
        try:
            import shutil
            if os.path.exists(self.unzipped_dir):
                shutil.rmtree(self.unzipped_dir)
            print(f"🧹 Cleaned up temporary files")
        except Exception as e:
            print(f"⚠️  Could not clean up temporary files: {e}")

# Main execution function
def run_enhanced_extraction(pptx_path=None, user_login="Rateb1223"):
    """
    Run enhanced 3D PowerPoint extraction
    
    Args:
        pptx_path (str): Path to PowerPoint file (auto-detected if None)
        user_login (str): User login for tracking
    
    Returns:
        str: Output directory path
    """
    extractor = Enhanced3DPowerPointExtractor(pptx_path, user_login)
    return extractor.extract_all()

# Enhanced demo execution
if __name__ == "__main__":
    # Current context
    current_time = datetime(2025, 6, 3, 16, 16, 32)
    user_login = "Rateb1223"
    
    print(f"🚀 Enhanced 3D PowerPoint Extractor Demo")
    print(f"👤 User: {user_login}")
    print(f"📅 Current Time: {current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC")
    print(f"📍 Working Directory: {os.getcwd()}")
    
    # Look for PowerPoint files
    pptx_files = [f for f in os.listdir('.') if f.endswith(('.pptx', '.ppt', '.potx'))]
    
    if not pptx_files:
        print("\n❌ No PowerPoint files found in current directory")
        print("📝 Please add PowerPoint files to test the enhanced extractor")
        print("\n📋 Features of this enhanced extractor:")
        print("  ✅ Advanced 3D chart detection and analysis")
        print("  ✅ XML-based 3D properties extraction")
        print("  ✅ Comprehensive Excel data matching")
        print("  ✅ Enhanced visualizations with 3D support")
        print("  ✅ Structured output with analysis reports")
        print("  ✅ Error handling and logging")
        print("  ✅ Dynamic paths and user tracking")
    else:
        print(f"\n🔍 Found {len(pptx_files)} PowerPoint files:")
        for i, file in enumerate(pptx_files, 1):
            print(f"  {i}. {file}")
        
        # Process all files or first file
        if len(pptx_files) == 1:
            selected_file = pptx_files[0]
        else:
            # For demo, process first file
            selected_file = pptx_files[0]
            print(f"\n📄 Processing first file: {selected_file}")
        
        # Run enhanced extraction
        result = run_enhanced_extraction(selected_file, user_login)
        
        if result:
            print(f"\n🎉 Enhanced extraction completed successfully!")
            print(f"📁 Check the output directory: {result}")
            print(f"\n📋 Generated content:")
            print(f"  - 3D chart visualizations and data")
            print(f"  - Excel data extraction and analysis")
            print(f"  - XML-based 3D properties analysis")
            print(f"  - Comprehensive analysis reports")
            print(f"  - Structured slide text extraction")
            print(f"  - Error logs and debugging information")
        else:
            print(f"\n❌ Enhanced extraction failed!")