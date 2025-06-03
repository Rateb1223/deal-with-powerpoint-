import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
import xml.etree.ElementTree as ET
from PIL import Image, ImageDraw, ImageFont
import io
import zipfile
import re
import base64
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
import json

class PowerPointToMarkdownConverter:
    def __init__(self, pptx_file_path, output_dir="converted_presentation"):
        self.pptx_file_path = pptx_file_path
        self.output_dir = output_dir
        self.images_dir = os.path.join(output_dir, "images")
        self.charts_dir = os.path.join(output_dir, "charts")
        self.data_dir = os.path.join(output_dir, "data")
        
        # Create directories
        for directory in [self.output_dir, self.images_dir, self.charts_dir, self.data_dir]:
            os.makedirs(directory, exist_ok=True)
        
        self.presentation = None
        self.chart_counter = 0
        self.image_counter = 0
        
    def convert(self):
        """Main conversion method"""
        try:
            self.presentation = Presentation(self.pptx_file_path)
            print(f"Converting presentation with {len(self.presentation.slides)} slides...")
            
            markdown_content = self._generate_markdown()
            
            # Save main markdown file
            markdown_file = os.path.join(self.output_dir, "presentation.md")
            with open(markdown_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            # Generate additional files
            self._generate_chart_summary()
            self._generate_metadata()
            
            print(f"✅ Conversion complete! Files saved to: {self.output_dir}")
            print(f"📄 Main file: {markdown_file}")
            
            return markdown_file
            
        except Exception as e:
            print(f"❌ Error during conversion: {e}")
            return None
    
    def _generate_markdown(self):
        """Generate the main markdown content"""
        markdown_lines = []
        
        # Header
        presentation_title = self._extract_presentation_title()
        markdown_lines.append(f"# {presentation_title}\n")
        markdown_lines.append(f"*Converted from PowerPoint on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n")
        markdown_lines.append("---\n")
        
        # Table of contents
        markdown_lines.append("## Table of Contents\n")
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_title = self._extract_slide_title(slide, i)
            markdown_lines.append(f"{i}. [{slide_title}](#slide-{i})\n")
        markdown_lines.append("\n---\n")
        
        # Process each slide
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            slide_markdown = self._process_slide(slide, slide_num)
            markdown_lines.extend(slide_markdown)
            markdown_lines.append("\n---\n")
        
        return "\n".join(markdown_lines)
    
    def _extract_presentation_title(self):
        """Extract presentation title from first slide or filename"""
        if self.presentation.slides:
            first_slide = self.presentation.slides[0]
            title = self._extract_slide_title(first_slide, 1)
            if title != "Slide 1":
                return title
        
        # Fallback to filename
        filename = os.path.basename(self.pptx_file_path)
        return os.path.splitext(filename)[0]
    
    def _extract_slide_title(self, slide, slide_num):
        """Extract title from a slide"""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except:
            pass
        
        # Look for the first text box that might be a title
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:  # Reasonable title length
                    return text
        
        return f"Slide {slide_num}"
    
    def _process_slide(self, slide, slide_num):
        """Process a single slide and return markdown lines"""
        markdown_lines = []
        
        # Slide header
        slide_title = self._extract_slide_title(slide, slide_num)
        markdown_lines.append(f"## Slide {slide_num}: {slide_title} {{#slide-{slide_num}}}\n")
        
        # Process slide layout and notes
        slide_layout = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown Layout"
        markdown_lines.append(f"*Layout: {slide_layout}*\n")
        
        # Extract slide notes if available
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_text = self._extract_notes(slide.notes_slide)
            if notes_text:
                markdown_lines.append(f"**Presenter Notes:** {notes_text}\n")
        
        # Process shapes in order of appearance
        text_content = []
        media_content = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_content = self._process_shape(shape, slide_num, shape_idx)
            if shape_content:
                if shape_content['type'] in ['text', 'title']:
                    text_content.append(shape_content)
                else:
                    media_content.append(shape_content)
        
        # Add text content
        for content in text_content:
            if content['content']:
                markdown_lines.extend(content['content'])
        
        # Add media content
        if media_content:
            markdown_lines.append("\n### Charts and Media\n")
            for content in media_content:
                if content['content']:
                    markdown_lines.extend(content['content'])
        
        return markdown_lines
    
    def _process_shape(self, shape, slide_num, shape_idx):
        """Process individual shapes and return content"""
        try:
            # Handle different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return self._process_chart(shape, slide_num, shape_idx)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._process_image(shape, slide_num, shape_idx)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return self._process_table(shape, slide_num, shape_idx)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
                return self._process_text(shape, slide_num, shape_idx)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                return self._process_group(shape, slide_num, shape_idx)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                return self._process_ole_object(shape, slide_num, shape_idx)
            
            else:
                # Handle other shape types
                return self._process_generic_shape(shape, slide_num, shape_idx)
                
        except Exception as e:
            print(f"⚠️  Error processing shape in slide {slide_num}: {e}")
            return None
    
    def _process_chart(self, chart_shape, slide_num, shape_idx):
        """Process charts and convert to images with data extraction"""
        try:
            self.chart_counter += 1
            chart_filename = f"chart_{slide_num}_{self.chart_counter}"
            
            # Extract chart data
            chart_data = self._extract_chart_data(chart_shape)
            
            # Create chart visualization
            chart_image_path = self._create_chart_image(chart_data, chart_filename)
            
            # Save chart data as JSON
            data_file = os.path.join(self.data_dir, f"{chart_filename}.json")
            with open(data_file, 'w', encoding='utf-8') as f:
                json.dump(chart_data, f, indent=2, ensure_ascii=False)
            
            # Create markdown content
            content = []
            content.append(f"#### Chart: {chart_data.get('title', 'Untitled Chart')}\n")
            
            if chart_image_path:
                rel_image_path = os.path.relpath(chart_image_path, self.output_dir)
                content.append(f"![{chart_data.get('title', 'Chart')}]({rel_image_path})\n")
            
            # Add chart details
            content.append(f"**Chart Type:** {chart_data.get('chart_type', 'Unknown')}\n")
            
            if chart_data.get('categories'):
                content.append(f"**Categories:** {', '.join(map(str, chart_data['categories'][:5]))}")
                if len(chart_data['categories']) > 5:
                    content.append("...")
                content.append("\n")
            
            # Add series information
            if chart_data.get('series_data'):
                content.append("**Data Series:**\n")
                for series in chart_data['series_data'][:3]:  # Limit to first 3 series
                    series_name = series.get('name', 'Unnamed Series')
                    values_count = len(series.get('values', []))
                    content.append(f"- {series_name} ({values_count} values)\n")
                
                if len(chart_data['series_data']) > 3:
                    content.append(f"- ... and {len(chart_data['series_data']) - 3} more series\n")
            
            # Link to raw data
            rel_data_path = os.path.relpath(data_file, self.output_dir)
            content.append(f"📊 [View Chart Data]({rel_data_path})\n")
            
            return {
                'type': 'chart',
                'content': content
            }
            
        except Exception as e:
            print(f"⚠️  Error processing chart in slide {slide_num}: {e}")
            return None
    
    def _extract_chart_data(self, chart_shape):
        """Extract comprehensive chart data"""
        try:
            chart = chart_shape.chart
            chart_data = {
                'title': 'Untitled Chart',
                'chart_type': str(chart.chart_type) if hasattr(chart, 'chart_type') else 'Unknown',
                'categories': [],
                'series_data': [],
                'has_3d': False
            }
            
            # Extract title
            try:
                if chart.has_title and chart.chart_title:
                    chart_data['title'] = chart.chart_title.text_frame.text
            except:
                pass
            
            # Check for 3D charts
            chart_type_str = str(chart.chart_type).lower()
            chart_data['has_3d'] = any(keyword in chart_type_str for keyword in ['3d', 'three', 'dimensional'])
            
            # Extract data from plots
            if hasattr(chart, 'plots') and chart.plots:
                for plot in chart.plots:
                    # Extract categories
                    if hasattr(plot, 'categories') and plot.categories:
                        try:
                            categories = []
                            for cat in plot.categories:
                                if hasattr(cat, 'label'):
                                    categories.append(str(cat.label))
                            chart_data['categories'] = categories
                        except:
                            pass
                    
                    # Extract series
                    if hasattr(plot, 'series'):
                        for series in plot.series:
                            try:
                                series_info = {
                                    'name': getattr(series, 'name', f'Series {len(chart_data["series_data"]) + 1}'),
                                    'values': []
                                }
                                
                                # Extract values
                                if hasattr(series, 'values') and series.values is not None:
                                    try:
                                        series_info['values'] = [float(v) if v is not None else 0 for v in series.values]
                                    except:
                                        series_info['values'] = list(series.values)
                                
                                chart_data['series_data'].append(series_info)
                            except:
                                continue
            
            return chart_data
            
        except Exception as e:
            print(f"Error extracting chart data: {e}")
            return {'title': 'Chart', 'chart_type': 'Unknown', 'categories': [], 'series_data': []}
    
    def _create_chart_image(self, chart_data, filename):
        """Create a visual representation of the chart"""
        try:
            plt.style.use('default')
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Determine chart type and create appropriate visualization
            chart_type = chart_data.get('chart_type', '').lower()
            categories = chart_data.get('categories', [])
            series_data = chart_data.get('series_data', [])
            
            if not series_data:
                # Create a placeholder chart
                ax.text(0.5, 0.5, f"Chart: {chart_data.get('title', 'No Data Available')}", 
                       ha='center', va='center', fontsize=12, transform=ax.transAxes)
                ax.set_xlim(0, 1)
                ax.set_ylim(0, 1)
            else:
                # Create chart based on type
                if 'bar' in chart_type or 'column' in chart_type:
                    self._create_bar_chart(ax, categories, series_data, chart_data.get('has_3d', False))
                elif 'line' in chart_type:
                    self._create_line_chart(ax, categories, series_data)
                elif 'pie' in chart_type:
                    self._create_pie_chart(ax, categories, series_data)
                elif 'scatter' in chart_type:
                    self._create_scatter_chart(ax, series_data)
                else:
                    # Default to bar chart
                    self._create_bar_chart(ax, categories, series_data, chart_data.get('has_3d', False))
            
            # Set title and labels
            title = chart_data.get('title', 'Chart')
            ax.set_title(title, fontsize=14, fontweight='bold', pad=20)
            
            # Add 3D indicator if applicable
            if chart_data.get('has_3d'):
                ax.text(0.02, 0.98, '3D Chart', transform=ax.transAxes, 
                       bbox=dict(boxstyle="round,pad=0.3", facecolor="lightblue"),
                       verticalalignment='top', fontsize=8)
            
            plt.tight_layout()
            
            # Save the chart
            chart_path = os.path.join(self.charts_dir, f"{filename}.png")
            plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close()
            
            return chart_path
            
        except Exception as e:
            print(f"Error creating chart image: {e}")
            return None
    
    def _create_bar_chart(self, ax, categories, series_data, is_3d=False):
        """Create a bar chart"""
        if not categories:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = np.arange(len(categories))
        width = 0.8 / len(series_data) if len(series_data) > 1 else 0.6
        
        colors = plt.cm.Set3(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                offset = (i - len(series_data)/2 + 0.5) * width
                bars = ax.bar(x + offset, values, width, label=series.get('name', f'Series {i+1}'), 
                             color=colors[i], alpha=0.8)
                
                # Add 3D effect if needed
                if is_3d:
                    for bar in bars:
                        bar.set_edgecolor('black')
                        bar.set_linewidth(0.5)
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        
        if len(series_data) > 1:
            ax.legend()
    
    def _create_line_chart(self, ax, categories, series_data):
        """Create a line chart"""
        if not categories:
            categories = [f"Point {i+1}" for i in range(len(series_data[0].get('values', [])))]
        
        x = range(len(categories))
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                ax.plot(x, values, marker='o', label=series.get('name', f'Series {i+1}'), 
                       color=colors[i], linewidth=2, markersize=4)
        
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.grid(True, alpha=0.3)
        
        if len(series_data) > 1:
            ax.legend()
    
    def _create_pie_chart(self, ax, categories, series_data):
        """Create a pie chart"""
        if series_data and series_data[0].get('values'):
            values = series_data[0]['values']
            labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
            
            # Filter out zero values
            filtered_data = [(label, value) for label, value in zip(labels, values) if value > 0]
            if filtered_data:
                labels, values = zip(*filtered_data)
                
                colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
                wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', 
                                                 colors=colors, startangle=90)
                
                # Improve text readability
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontweight('bold')
        
        ax.axis('equal')
    
    def _create_scatter_chart(self, ax, series_data):
        """Create a scatter chart"""
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        
        for i, series in enumerate(series_data[:2]):  # Use first two series for x and y
            values = series.get('values', [])
            if values:
                if i == 0:
                    x_values = values
                    x_label = series.get('name', 'X Values')
                else:
                    y_values = values[:len(x_values)]  # Match lengths
                    y_label = series.get('name', 'Y Values')
                    ax.scatter(x_values, y_values, color=colors[0], alpha=0.6, s=50)
                    ax.set_xlabel(x_label)
                    ax.set_ylabel(y_label)
                    break
        
        ax.grid(True, alpha=0.3)
    
    def _process_image(self, image_shape, slide_num, shape_idx):
        """Process images and save them"""
        try:
            self.image_counter += 1
            
            # Extract image
            image = image_shape.image
            image_ext = getattr(image, 'ext', 'png')
            if image_ext.startswith('.'):
                image_ext = image_ext[1:]
            
            # Save image
            image_filename = f"image_{slide_num}_{self.image_counter}.{image_ext}"
            image_path = os.path.join(self.images_dir, image_filename)
            
            with open(image_path, 'wb') as f:
                f.write(image.blob)
            
            # Create markdown content
            rel_image_path = os.path.relpath(image_path, self.output_dir)
            content = [f"![Image]({rel_image_path})\n"]
            
            return {
                'type': 'image',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing image: {e}")
            return None
    
    def _process_table(self, table_shape, slide_num, shape_idx):
        """Process tables and convert to markdown tables"""
        try:
            table = table_shape.table
            markdown_table = []
            
            # Extract table data
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            if not table_data:
                return None
            
            # Create markdown table
            markdown_table.append("| " + " | ".join(table_data[0]) + " |")
            markdown_table.append("|" + "---|" * len(table_data[0]))
            
            for row in table_data[1:]:
                markdown_table.append("| " + " | ".join(row) + " |")
            
            markdown_table.append("")  # Empty line after table
            
            return {
                'type': 'table',
                'content': markdown_table
            }
            
        except Exception as e:
            print(f"Error processing table: {e}")
            return None
    
    def _process_text(self, text_shape, slide_num, shape_idx):
        """Process text content"""
        try:
            if not hasattr(text_shape, 'text_frame') or not text_shape.text_frame:
                return None
            
            text_content = []
            text_frame = text_shape.text_frame
            
            # Check if this is a title
            is_title = (hasattr(text_shape, 'placeholder_format') and 
                       text_shape.placeholder_format and 
                       'title' in str(text_shape.placeholder_format.type).lower())
            
            for paragraph in text_frame.paragraphs:
                if paragraph.text.strip():
                    # Determine formatting level
                    level = paragraph.level if hasattr(paragraph, 'level') else 0
                    text = paragraph.text.strip()
                    
                    if is_title and level == 0:
                        text_content.append(f"### {text}\n")
                    elif level == 0:
                        text_content.append(f"{text}\n")
                    else:
                        bullet = "  " * level + "- "
                        text_content.append(f"{bullet}{text}\n")
            
            if text_content:
                text_content.append("")  # Empty line
            
            return {
                'type': 'title' if is_title else 'text',
                'content': text_content
            }
            
        except Exception as e:
            print(f"Error processing text: {e}")
            return None
    
    def _process_group(self, group_shape, slide_num, shape_idx):
        """Process grouped shapes"""
        try:
            group_content = []
            
            for i, shape in enumerate(group_shape.shapes):
                shape_content = self._process_shape(shape, slide_num, f"{shape_idx}_g{i}")
                if shape_content and shape_content.get('content'):
                    group_content.extend(shape_content['content'])
            
            if group_content:
                return {
                    'type': 'group',
                    'content': group_content
                }
            
        except Exception as e:
            print(f"Error processing group: {e}")
        
        return None
    
    def _process_ole_object(self, ole_shape, slide_num, shape_idx):
        """Process embedded OLE objects"""
        try:
            prog_id = getattr(ole_shape.ole_format, 'prog_id', 'Unknown')
            
            content = []
            content.append(f"**Embedded Object:** {prog_id}\n")
            
            # Try to extract the object
            if hasattr(ole_shape.ole_format, 'blob'):
                # Determine file extension
                if 'excel' in prog_id.lower():
                    ext = 'xlsx'
                elif 'word' in prog_id.lower():
                    ext = 'docx'
                else:
                    ext = 'bin'
                
                # Save the embedded object
                obj_filename = f"embedded_{slide_num}_{shape_idx}.{ext}"
                obj_path = os.path.join(self.data_dir, obj_filename)
                
                with open(obj_path, 'wb') as f:
                    f.write(ole_shape.ole_format.blob)
                
                rel_obj_path = os.path.relpath(obj_path, self.output_dir)
                content.append(f"📎 [Download Embedded Object]({rel_obj_path})\n")
            
            return {
                'type': 'ole_object',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing OLE object: {e}")
            return None
    
    def _process_generic_shape(self, shape, slide_num, shape_idx):
        """Process other types of shapes"""
        try:
            shape_type = str(shape.shape_type)
            content = [f"*Shape: {shape_type}*\n"]
            
            # Try to extract any text
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    content.append(f"{text}\n")
            
            return {
                'type': 'generic_shape',
                'content': content
            }
            
        except Exception as e:
            print(f"Error processing generic shape: {e}")
            return None
    
    def _extract_notes(self, notes_slide):
        """Extract presenter notes from slide"""
        try:
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        except:
            pass
        return ""
    
    def _generate_chart_summary(self):
        """Generate a summary of all charts"""
        summary_file = os.path.join(self.output_dir, "charts_summary.md")
        
        with open(summary_file, 'w', encoding='utf-8') as f:
            f.write("# Charts Summary\n\n")
            f.write(f"Total charts extracted: {self.chart_counter}\n\n")
            
            # List all chart files
            chart_files = [f for f in os.listdir(self.charts_dir) if f.endswith('.png')]
            data_files = [f for f in os.listdir(self.data_dir) if f.endswith('.json')]
            
            f.write("## Chart Images\n\n")
            for chart_file in sorted(chart_files):
                f.write(f"- ![{chart_file}](charts/{chart_file})\n")
            
            f.write("\n## Chart Data Files\n\n")
            for data_file in sorted(data_files):
                f.write(f"- [📊 {data_file}](data/{data_file})\n")
    
    def _generate_metadata(self):
        """Generate metadata about the conversion"""
        metadata = {
            'source_file': os.path.basename(self.pptx_file_path),
            'conversion_date': datetime.now().isoformat(),
            'total_slides': len(self.presentation.slides) if self.presentation else 0,
            'charts_extracted': self.chart_counter,
            'images_extracted': self.image_counter,
            'output_structure': {
                'main_file': 'presentation.md',
                'charts_directory': 'charts/',
                'images_directory': 'images/',
                'data_directory': 'data/',
                'summary_file': 'charts_summary.md'
            }
        }
        
        metadata_file = os.path.join(self.output_dir, "metadata.json")
        with open(metadata_file, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)

def convert_powerpoint_to_markdown(pptx_file, output_dir=None):
    """
    Main function to convert PowerPoint to Markdown
    
    Args:
        pptx_file (str): Path to PowerPoint file
        output_dir (str): Output directory (optional)
    
    Returns:
        str: Path to the generated markdown file
    """
    if output_dir is None:
        base_name = os.path.splitext(os.path.basename(pptx_file))[0]
        output_dir = f"{base_name}_markdown"
    
    converter = PowerPointToMarkdownConverter(pptx_file, output_dir)
    return converter.convert()

# Example usage
if __name__ == "__main__":
    # Replace with your PowerPoint file path
    pptx_file = "test.pptx"
    
    if os.path.exists(pptx_file):
        print("🚀 Starting PowerPoint to Markdown conversion...")
        result = convert_powerpoint_to_markdown(pptx_file)
        
        if result:
            print(f"✅ Conversion successful!")
            print(f"📁 Output directory: {os.path.dirname(result)}")
        else:
            print("❌ Conversion failed!")