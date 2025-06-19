import os
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook
import pandas as pd
import json
import xml.etree.ElementTree as ET
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
import numpy as np
from datetime import datetime
import glob
import re
from pathlib import Path
import shutil

# --- Constants and Utility Functions (if any, could be moved to a separate utils.py later) ---
# For now, keeping them within this file for self-containment as requested.

# --- 1. PowerPointFileHandler Class ---
class PowerPointFileHandler:
    """
    Handles file-system related operations for PowerPoint extraction.
    This includes finding the PowerPoint file, creating and managing
    output directories, unzipping the PPTX, and cleaning up temporary files.
    """
    def __init__(self, pptx_path: str = None, user_login: str = "default_user", current_time: datetime = None):
        """
        Initializes the file handler with paths and user/time metadata.

        Args:
            pptx_path (str): Path to the PowerPoint file. If None, it attempts to auto-detect.
            user_login (str): User identifier for output directory naming.
            current_time (datetime): Timestamp for output directory naming. If None, uses current UTC time.
        """
        self.user_login = user_login
        self.current_time = current_time if current_time else datetime.utcnow() # Ensure current_time is set
        self.pptx_path = pptx_path if pptx_path else self._find_powerpoint_file()

        # Validate pptx_path
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {self.pptx_path}")

        # Dynamic output directory with timestamp and user
        base_name = Path(self.pptx_path).stem
        timestamp = self.current_time.strftime("%Y%m%d_%H%M%S")
        self.output_dir = f"enhanced_extraction_{base_name}_{self.user_login}_{timestamp}"
        self.unzipped_dir = f"temp_unzipped_{timestamp}"

        # Create structured output directories
        self.dirs = {
            'main': self.output_dir,
            'charts_3d': os.path.join(self.output_dir, "charts_3d"),
            'charts_2d': os.path.join(self.output_dir, "charts_2d"),
            'tables': os.path.join(self.output_dir, "tables"),
            'excel_data': os.path.join(self.output_dir, "excel_data"),
            'xml_data': os.path.join(self.output_dir, "xml_analysis"),
            'slide_text': os.path.join(self.output_dir, "slide_text"),
            'errors': os.path.join(self.output_dir, "errors"),
            'analysis': os.path.join(self.output_dir, "analysis"),
            'markdown_slides': os.path.join(self.output_dir, "markdown_slides") # Add markdown dir here
        }

        for dir_path in self.dirs.values():
            os.makedirs(dir_path, exist_ok=True)

        print(f"🚀 PowerPoint File Handler Initialized")
        print(f"📄 Processing: {os.path.basename(self.pptx_path)}")
        print(f"📂 Output Directory: {self.output_dir}")

    def _find_powerpoint_file(self) -> str:
        """
        Auto-detects a PowerPoint file in the current working directory.
        Prioritizes .pptx files.

        Returns:
            str: The path to the detected PowerPoint file, or None if none found.
        """
        patterns = ["*.pptx", "*.ppt", "*.potx"]
        files = []

        for pattern in patterns:
            files.extend(glob.glob(pattern))

        if files:
            pptx_files = [f for f in files if f.endswith('.pptx')]
            selected_file = pptx_files[0] if pptx_files else files[0]
            print(f"🔍 Auto-detected PowerPoint file: {selected_file}")
            return selected_file
        else:
            print("❌ No PowerPoint files found in current directory")
            return None

    def unzip_pptx(self):
        """
        Unzips the PowerPoint file to a temporary directory.
        """
        print(f"\n📦 Unzipping PowerPoint file to: {self.unzipped_dir}")
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                zip_ref.extractall(self.unzipped_dir)
            print("✅ PowerPoint file unzipped successfully.")
        except Exception as e:
            print(f"❌ Error unzipping PowerPoint file: {e}")
            raise

    def cleanup_temp_files(self):
        """
        Removes the temporary unzipped directory.
        """
        print(f"\n🧹 Cleaning up temporary files in: {self.unzipped_dir}")
        try:
            if os.path.exists(self.unzipped_dir):
                shutil.rmtree(self.unzipped_dir)
            print("✅ Temporary files cleaned up.")
        except Exception as e:
            print(f"⚠️  Could not clean up temporary files: {e}")

    def get_paths(self) -> dict:
        """
        Returns a dictionary of all managed directory paths.
        """
        return self.dirs

    def get_pptx_path(self) -> str:
        """
        Returns the path to the PowerPoint file being processed.
        """
        return self.pptx_path

    def get_unzipped_dir(self) -> str:
        """
        Returns the path to the temporary unzipped directory.
        """
        return self.unzipped_dir

# --- 2. ExcelDataExtractor Class ---
class ExcelDataExtractor:
    """
    Extracts embedded Excel data from the unzipped PowerPoint file
    and performs initial analysis on the sheets.
    """
    def __init__(self, unzipped_dir: str, output_dirs: dict, user_login: str, current_time: datetime):
        """
        Initializes the Excel data extractor.

        Args:
            unzipped_dir (str): Path to the temporary unzipped PowerPoint directory.
            output_dirs (dict): Dictionary of output directories for saving data.
            user_login (str): User identifier for metadata.
            current_time (datetime): Timestamp for metadata.
        """
        self.unzipped_dir = unzipped_dir
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time
        self.excel_counter = 0
        self.embedded_data = [] # Stores extracted excel data
        self.errors = [] # Stores errors related to excel extraction

    def extract_embedded_excel_data(self):
        """
        Extracts embedded Excel data from the unzipped PowerPoint.
        Loads each worksheet, extracts its content, analyzes its structure,
        and saves it as CSV and JSON.
        """
        print(f"\n📊 Extracting embedded Excel data...")

        embedded_excel_dir = os.path.join(self.unzipped_dir, "ppt", "embeddings")

        if not os.path.exists(embedded_excel_dir):
            print("⚠️  No embedded Excel files found.")
            return

        embedded_excels = sorted([f for f in os.listdir(embedded_excel_dir)
                                if f.endswith((".xlsx", ".xls"))])

        print(f"📁 Found {len(embedded_excels)} embedded Excel files.")

        for file_idx, file in enumerate(embedded_excels):
            try:
                self.excel_counter += 1
                file_path = os.path.join(embedded_excel_dir, file)

                print(f"  📄 Processing Excel {self.excel_counter}: {file}")

                wb = load_workbook(file_path, data_only=True)

                excel_data = {
                    'file_name': file,
                    'file_index': self.excel_counter,
                    'extraction_metadata': {
                        'extracted_by': self.user_login,
                        'extraction_time': self.current_time.isoformat(),
                        'file_size': os.path.getsize(file_path)
                    },
                    'sheets': []
                }

                for sheet_idx, sheet in enumerate(wb.worksheets):
                    print(f"    📋 Sheet: {sheet.title}")
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell).strip() if cell is not None else "" for cell in row]
                        rows.append(row_data)

                    sheet_analysis = self._analyze_excel_sheet(rows, sheet.title)

                    sheet_data = {
                        'sheet_name': sheet.title,
                        'sheet_index': sheet_idx,
                        'rows': len(rows),
                        'columns': len(rows[0]) if rows else 0,
                        'data': rows,
                        'analysis': sheet_analysis
                    }

                    excel_data['sheets'].append(sheet_data)

                    if rows:
                        csv_path = os.path.join(self.output_dirs['excel_data'],
                                              f"excel_{self.excel_counter}_sheet_{sheet_idx}_{sheet.title}.csv")
                        df = pd.DataFrame(rows[1:], columns=rows[0] if rows else None)
                        df.to_csv(csv_path, index=False, encoding='utf-8')

                json_path = os.path.join(self.output_dirs['excel_data'], f"excel_{self.excel_counter}_metadata.json")
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(excel_data, f, indent=2, ensure_ascii=False)

                self.embedded_data.append(excel_data)

            except Exception as e:
                error_info = {
                    'file': file,
                    'error': str(e),
                    'error_type': 'excel_extraction'
                }
                self.errors.append(error_info)
                print(f"    ⚠️  Error processing Excel file: {e}")

    def _analyze_excel_sheet(self, rows: list, sheet_name: str) -> dict:
        """
        Analyzes Excel sheet structure for better chart matching.
        Identifies headers, numeric/text columns, and potential categories/values.

        Args:
            rows (list): List of rows from the Excel sheet.
            sheet_name (str): Name of the sheet being analyzed.

        Returns:
            dict: Analysis results including column types and data structure.
        """
        analysis = {
            'has_headers': False,
            'numeric_columns': [],
            'text_columns': [],
            'potential_categories': [],
            'potential_values': [],
            'data_structure': 'unknown'
        }

        if not rows or len(rows) < 2:
            return analysis

        headers = rows[0]
        data_rows = rows[1:]

        if headers and any(isinstance(h, str) and h.strip() for h in headers):
            analysis['has_headers'] = True

        for col_idx, header in enumerate(headers):
            column_values = [row[col_idx] if col_idx < len(row) else '' for row in data_rows]

            numeric_count = 0
            for val in column_values:
                try:
                    float(str(val).replace(',', ''))
                    numeric_count += 1
                except ValueError:
                    pass

            if numeric_count > len(column_values) * 0.7:
                analysis['numeric_columns'].append(col_idx)
                analysis['potential_values'].extend(column_values)
            else:
                analysis['text_columns'].append(col_idx)
                if col_idx == 0:
                    analysis['potential_categories'] = column_values

        if len(analysis['text_columns']) >= 1 and len(analysis['numeric_columns']) >= 1:
            analysis['data_structure'] = 'chart_suitable'
        elif len(analysis['numeric_columns']) > 1:
            analysis['data_structure'] = 'multi_series'
        else:
            analysis['data_structure'] = 'simple_table'

        return analysis

    def get_extracted_data(self) -> list:
        """
        Returns the list of all extracted embedded Excel data.
        """
        return self.embedded_data

    def get_excel_count(self) -> int:
        """
        Returns the total count of extracted Excel files.
        """
        return self.excel_counter

    def get_errors(self) -> list:
        """
        Returns any errors encountered during Excel extraction.
        """
        return self.errors

# --- 3. ChartXmlAnalyzer Class ---
class ChartXmlAnalyzer:
    """
    Focuses solely on parsing chart XML files to detect 3D properties
    and other chart metadata.
    """
    def __init__(self, unzipped_dir: str, output_dirs: dict, user_login: str, current_time: datetime):
        """
        Initializes the XML chart analyzer.

        Args:
            unzipped_dir (str): Path to the temporary unzipped PowerPoint directory.
            output_dirs (dict): Dictionary of output directories for saving analysis.
            user_login (str): User identifier for metadata.
            current_time (datetime): Timestamp for metadata.
        """
        self.unzipped_dir = unzipped_dir
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time
        self.xml_charts = {} # Stores analyzed XML chart data

    def extract_xml_chart_data(self):
        """
        Finds and processes chart XML files within the unzipped PowerPoint,
        analyzing each for 3D properties and other relevant chart information.
        """
        print(f"\n🔍 Extracting XML chart data for 3D analysis...")

        charts_dir = os.path.join(self.unzipped_dir, "ppt", "charts")

        if not os.path.exists(charts_dir):
            print("⚠️  No chart XML files found.")
            return

        chart_files = [f for f in os.listdir(charts_dir) if f.endswith('.xml')]
        print(f"📊 Found {len(chart_files)} chart XML files.")

        for chart_file in chart_files:
            try:
                chart_path = os.path.join(charts_dir, chart_file)

                with open(chart_path, 'r', encoding='utf-8') as f:
                    xml_content = f.read()

                chart_analysis = self._analyze_chart_xml(xml_content, chart_file)

                chart_id = chart_file.replace('.xml', '')
                self.xml_charts[chart_id] = chart_analysis

                analysis_path = os.path.join(self.output_dirs['xml_data'], f"{chart_id}_analysis.json")
                with open(analysis_path, 'w', encoding='utf-8') as f:
                    json.dump(chart_analysis, f, indent=2, ensure_ascii=False)

                print(f"  📄 Analyzed: {chart_file} (3D: {chart_analysis.get('is_3d', False)})")

            except Exception as e:
                print(f"    ⚠️  Error analyzing {chart_file}: {e}")

    def _analyze_chart_xml(self, xml_content: str, filename: str) -> dict:
        """
        Parses the XML content of a chart to detect if it's 3D,
        its type, and extracts 3D view properties.

        Args:
            xml_content (str): The XML content of the chart.
            filename (str): The original filename of the XML.

        Returns:
            dict: Analysis results including 3D detection and properties.
        """
        analysis = {
            'filename': filename,
            'is_3d': False,
            'chart_type': 'unknown',
            '3d_properties': {},
            'series_count': 0,
            'categories_count': 0,
            'extraction_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            }
        }

        try:
            root = ET.fromstring(xml_content)
            namespace = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
            if '}' in root.tag:
                ns_uri = root.tag.split('}')[0].strip('{')
                namespace = {'c': ns_uri}

            chart_type_elem = root.find('.//c:chart', namespace)
            if chart_type_elem is not None:
                plot_area = chart_type_elem.find('.//c:plotArea', namespace)
                if plot_area is not None:
                    chart_types = [
                        'bar3DChart', 'line3DChart', 'pie3DChart', 'area3DChart',
                        'surface3DChart', 'bubbleChart', 'barChart', 'lineChart'
                    ]

                    for chart_type in chart_types:
                        type_elem = plot_area.find(f'.//c:{chart_type}', namespace)
                        if type_elem is not None:
                            analysis['chart_type'] = chart_type
                            analysis['is_3d'] = '3d' in chart_type.lower() or 'surface' in chart_type.lower()
                            break

            if analysis['is_3d']:
                view3d = root.find('.//c:view3D', namespace)
                if view3d is not None:
                    properties = {}
                    for prop in ['rotX', 'rotY', 'depthPercent', 'rAngAx', 'perspective']:
                        elem = view3d.find(f'./c:{prop}', namespace)
                        if elem is not None and 'val' in elem.attrib:
                            properties[prop] = elem.attrib['val']
                    analysis['3d_properties'] = properties

            series_elems = root.findall('.//c:ser', namespace)
            analysis['series_count'] = len(series_elems)

            cat_elems = root.findall('.//c:cat', namespace)
            analysis['categories_count'] = len(cat_elems)

        except Exception as e:
            analysis['xml_error'] = str(e)

        return analysis

    def get_xml_charts_analysis(self) -> dict:
        """
        Returns the dictionary of analyzed XML chart data.
        """
        return self.xml_charts

# --- 4. ChartVisualizer Class ---
class ChartVisualizer:
    """
    Encapsulates all chart visualization logic (both 2D and 3D).
    Uses Matplotlib to render charts as PNG images.
    """
    def __init__(self, output_dirs: dict, user_login: str, current_time: datetime):
        """
        Initializes the chart visualizer.

        Args:
            output_dirs (dict): Dictionary of output directories for saving images.
            user_login (str): User identifier for metadata on the chart image.
            current_time (datetime): Timestamp for metadata on the chart image.
        """
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time

    def create_chart_visualization(self, chart_data: dict, chart_idx: int, doc_id: str, slide_num: int):
        """
        Creates an enhanced chart visualization (2D or 3D) and saves it as a PNG.

        Args:
            chart_data (dict): Dictionary containing all extracted chart data.
            chart_idx (int): The index of the chart being visualized (for filename).
            doc_id (str): Document ID (for filename).
            slide_num (int): Slide number (for filename).
        """
        try:
            chart_id = chart_data['chart_id']
            is_3d = chart_data['is_3d']

            if is_3d:
                fig = plt.figure(figsize=(12, 9))
                ax = fig.add_subplot(111, projection='3d')
                self._create_3d_chart(ax, chart_data)
                self._apply_3d_properties(ax, chart_data.get('3d_properties', {}))
                chart_path = os.path.join(self.output_dirs['charts_3d'], f"{chart_id}.png")
            else:
                fig, ax = plt.subplots(figsize=(10, 6))
                self._create_2d_chart(ax, chart_data)
                chart_path = os.path.join(self.output_dirs['charts_2d'], f"{chart_id}.png")

            title = f"{'🎲 3D' if is_3d else '📊'} Chart: {chart_data['chart_type']}"
            ax.set_title(title, fontsize=14, fontweight='bold', pad=20)

            metadata_text = f"Extracted by {self.user_login} | {self.current_time.strftime('%Y-%m-%d %H:%M')}"
            if is_3d:
                ax.text2D(0.02, 0.02, metadata_text, transform=ax.transAxes, fontsize=8, alpha=0.7)
            else:
                ax.text(0.02, 0.02, metadata_text, transform=ax.transAxes, fontsize=8, alpha=0.7)

            plt.tight_layout()
            plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close()

            chart_data['visualization_path'] = chart_path
            print(f"      💾 Saved visualization: {os.path.basename(chart_path)}")

            # Save the image from the chart shape if available (this is specific to original shape object)
            # This part usually requires the actual chart_shape object which is not passed here.
            # If the original image extraction is desired, it needs to happen where chart_shape is accessible.
            # For now, we only save the *generated* visualization.

        except Exception as e:
            print(f"      ⚠️  Error creating visualization for {chart_id}: {e}")
            chart_data['visualization_path'] = None # Mark as failed if error occurs

    def _create_3d_chart(self, ax: Axes3D, chart_data: dict):
        """
        Creates the appropriate 3D chart visualization based on chart type.
        """
        categories = chart_data.get('categories', [])
        series_data = chart_data.get('series_data', [])
        chart_type = chart_data.get('chart_type', '').lower()

        if not series_data:
            ax.text2D(0.5, 0.5, "No Data Available", transform=ax.transAxes,
                     ha='center', va='center', fontsize=14)
            return

        if 'column' in chart_type or 'bar' in chart_type:
            self._create_3d_bar_chart(ax, categories, series_data)
        elif 'line' in chart_type:
            self._create_3d_line_chart(ax, categories, series_data)
        elif 'pie' in chart_type:
            self._create_3d_pie_chart(ax, categories, series_data)
        elif 'surface' in chart_type or 'area' in chart_type:
            self._create_3d_surface_chart(ax, categories, series_data)
        else:
            self._create_3d_bar_chart(ax, categories, series_data)

        ax.set_xlabel('Categories')
        ax.set_ylabel('Series')
        ax.set_zlabel('Values')

    def _create_3d_bar_chart(self, ax: Axes3D, categories: list, series_data: list):
        """Creates a 3D bar chart."""
        if not categories:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]

        colors = plt.cm.viridis(np.linspace(0, 1, len(series_data)))
        x_pos = np.arange(len(categories))
        y_pos = np.arange(len(series_data))
        width = 0.6
        depth = 0.6

        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                for j, val in enumerate(values):
                    if j < len(x_pos):
                        x = x_pos[j]
                        y = y_pos[i] if i < len(y_pos) else i
                        z = 0
                        dx = width
                        dy = depth
                        dz = max(float(val), 0.1) if isinstance(val, (int, float)) else 0.1
                        ax.bar3d(x, y, z, dx, dy, dz, color=colors[i], alpha=0.8)

        ax.set_xticks(x_pos + width/2)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.set_yticks(y_pos + depth/2)
        ax.set_yticklabels([s.get('name', f'Series {i+1}') for i, s in enumerate(series_data)])

    def _create_3d_line_chart(self, ax: Axes3D, categories: list, series_data: list):
        """Creates a 3D line chart."""
        colors = plt.cm.tab10(np.linspace(0, 1, len(series_data)))
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            z_values = series.get('z_values', [])
            if values:
                x = np.arange(len(values))
                y = z_values if z_values and len(z_values) == len(values) else [i] * len(values)
                z = [float(v) if isinstance(v, (int, float)) else 0 for v in values]
                ax.plot(x, y, z, color=colors[i], linewidth=3, marker='o', markersize=5,
                       label=series.get('name', f'Series {i+1}'))
        ax.legend()

    def _create_3d_pie_chart(self, ax: Axes3D, categories: list, series_data: list):
        """Creates a representation of a 3D pie chart."""
        if not series_data or not series_data[0].get('values'):
            return
        values = series_data[0]['values']
        labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
        angles = np.linspace(0, 2*np.pi, len(values), endpoint=False)
        colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
        for i, (angle, value, color) in enumerate(zip(angles, values, colors)):
            if isinstance(value, (int, float)) and value > 0:
                theta = np.linspace(angle, angle + 2*np.pi/len(values), 20)
                r = 3
                height = value / max(values) if max(values) > 0 else 0.1
                x = r * np.cos(theta)
                y = r * np.sin(theta)
                z = np.full_like(x, height)
                ax.plot(x, y, z, color=color, linewidth=3)

    def _create_3d_surface_chart(self, ax: Axes3D, categories: list, series_data: list):
        """Creates a 3D surface chart."""
        if not series_data:
            return
        x_size = len(categories) if categories else len(series_data[0].get('values', []))
        y_size = len(series_data)
        if x_size == 0 or y_size == 0:
            return
        x = np.arange(x_size)
        y = np.arange(y_size)
        X, Y = np.meshgrid(x, y)
        Z = np.zeros_like(X, dtype=float)
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            for j, val in enumerate(values):
                if j < Z.shape[1] and isinstance(val, (int, float)):
                    Z[i, j] = float(val)
        surf = ax.plot_surface(X, Y, Z, cmap='viridis', alpha=0.8)
        plt.colorbar(surf, ax=ax, shrink=0.5, aspect=5)

    def _create_2d_chart(self, ax: plt.Axes, chart_data: dict):
        """
        Creates the appropriate 2D chart visualization based on chart type.
        """
        categories = chart_data.get('categories', [])
        series_data = chart_data.get('series_data', [])
        chart_type = chart_data.get('chart_type', '').lower()

        if not series_data:
            ax.text(0.5, 0.5, "No Data Available", transform=ax.transAxes,
                   ha='center', va='center', fontsize=14)
            return

        if 'bar' in chart_type or 'column' in chart_type:
            self._create_2d_bar_chart(ax, categories, series_data)
        elif 'line' in chart_type:
            self._create_2d_line_chart(ax, categories, series_data)
        elif 'pie' in chart_type:
            self._create_2d_pie_chart(ax, categories, series_data)
        else:
            self._create_2d_bar_chart(ax, categories, series_data)

    def _create_2d_bar_chart(self, ax: plt.Axes, categories: list, series_data: list):
        """Creates a 2D bar chart."""
        if not categories:
            categories = [f"Item {i+1}" for i in range(len(series_data[0].get('values', [])))]
        x = np.arange(len(categories))
        width = 0.8 / len(series_data) if len(series_data) > 1 else 0.6
        colors = plt.cm.Set3(np.linspace(0, 1, len(series_data)))
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                offset = (i - len(series_data)/2 + 0.5) * width
                ax.bar(x + offset, values, width, label=series.get('name', f'Series {i+1}'),
                      color=colors[i], alpha=0.8)
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        if len(series_data) > 1:
            ax.legend()

    def _create_2d_line_chart(self, ax: plt.Axes, categories: list, series_data: list):
        """Creates a 2D line chart."""
        if not categories:
            categories = [f"Point {i+1}" for i in range(len(series_data[0].get('values', [])))]
        x = range(len(categories))
        colors = plt.cm.Set1(np.linspace(0, 1, len(series_data)))
        for i, series in enumerate(series_data):
            values = series.get('values', [])
            if values:
                ax.plot(x, values, marker='o', label=series.get('name', f'Series {i+1}'),
                       color=colors[i], linewidth=2, markersize=4)
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.grid(True, alpha=0.3)
        if len(series_data) > 1:
            ax.legend()

    def _create_2d_pie_chart(self, ax: plt.Axes, categories: list, series_data: list):
        """Creates a 2D pie chart."""
        if series_data and series_data[0].get('values'):
            values = series_data[0]['values']
            labels = categories if categories else [f"Slice {i+1}" for i in range(len(values))]
            filtered_data = [(label, float(value)) for label, value in zip(labels, values)
                           if isinstance(value, (int, float)) and value > 0]
            if filtered_data:
                labels, values = zip(*filtered_data)
                colors = plt.cm.Set3(np.linspace(0, 1, len(values)))
                wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%',
                                                 colors=colors, startangle=90)
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontweight('bold')
        ax.axis('equal')

    def _apply_3d_properties(self, ax: Axes3D, properties: dict):
        """
        Applies 3D view properties (rotation, perspective) to the Matplotlib 3D axes.
        """
        try:
            elev, azim = 20, -35 # Default view
            if 'rotX' in properties:
                elev = int(properties['rotX']) / 60000
            if 'rotY' in properties:
                azim = int(properties['rotY']) / 60000
            ax.view_init(elev=elev, azim=azim)
        except Exception:
            ax.view_init(elev=20, azim=-35) # Fallback to default

# --- 5. ChartDataProcessor Class ---
class ChartDataProcessor:
    """
    Handles the extraction of data from PowerPoint chart objects,
    matching them with extracted Excel data, determining 3D properties,
    and preparing data for visualization.
    """
    def __init__(self, embedded_excel_data: list, xml_charts_analysis: dict,
                 user_login: str, current_time: datetime):
        """
        Initializes the chart data processor.

        Args:
            embedded_excel_data (list): List of extracted embedded Excel data.
            xml_charts_analysis (dict): Dictionary of analyzed XML chart data.
            user_login (str): User identifier for metadata.
            current_time (datetime): Timestamp for metadata.
        """
        self.embedded_excel_data = embedded_excel_data
        self.xml_charts_analysis = xml_charts_analysis
        self.user_login = user_login
        self.current_time = current_time

    def process_chart_data(self, chart_shape, slide_num: int, chart_counter: int) -> dict:
        """
        Processes a single PowerPoint chart shape, extracting its data,
        matching with embedded Excel, detecting 3D properties, and creating
        a structured data dictionary for it.

        Args:
            chart_shape: The PowerPoint chart shape object.
            slide_num (int): The slide number where the chart is located.
            chart_counter (int): The current chart count (for unique IDs).

        Returns:
            dict: A dictionary containing all extracted and processed chart data, or None if error.
        """
        try:
            chart = chart_shape.chart
            chart_id = f"chart_{slide_num}_{chart_counter + 1}"

            chart_data = {
                'chart_id': chart_id,
                'slide_number': slide_num,
                'chart_type': str(chart.chart_type),
                'is_3d': False,
                '3d_properties': {},
                'excel_data': None,
                'categories': [],
                'series_data': [],
                'text_representation': [],
                'image_path': None, # To be filled by visualizer or image extraction
                'visualization_path': None, # To be filled by visualizer
                'extraction_metadata': {
                    'extracted_by': self.user_login,
                    'extraction_time': self.current_time.isoformat(),
                    'excel_index': chart_counter # This links to the excel data order
                }
            }

            # Enhanced 3D detection based on type string
            chart_type_str = str(chart.chart_type).lower()
            chart_data['is_3d'] = any(keyword in chart_type_str for keyword in
                                    ['3d', 'three', 'dimensional', 'surface', 'cone', 'pyramid'])

            # Get XML chart analysis if available and merge 3D properties
            xml_chart_id = f"chart{chart_counter + 1}" # XML charts are named chart1, chart2 etc.
            if xml_chart_id in self.xml_charts_analysis:
                xml_analysis = self.xml_charts_analysis[xml_chart_id]
                chart_data['is_3d'] = chart_data['is_3d'] or xml_analysis.get('is_3d', False)
                chart_data['3d_properties'].update(xml_analysis.get('3d_properties', {}))

            print(f"    {'🎲' if chart_data['is_3d'] else '📊'} Chart: {chart_id} ({'3D' if chart_data['is_3d'] else '2D'})")

            # Match with embedded Excel data
            if chart_counter < len(self.embedded_excel_data):
                excel_file = self.embedded_excel_data[chart_counter]
                if excel_file['sheets']:
                    excel_data = excel_file['sheets'][0] # Use first sheet
                    chart_data['excel_data'] = excel_data
                    self._extract_chart_data_from_excel(chart_data, excel_data)

            # Create text representation
            text_lines = [f"CHART START - {chart_data['chart_type']}"]
            if chart_data['is_3d']:
                text_lines.append("🎲 3D CHART DETECTED")
                if chart_data['3d_properties']:
                    text_lines.append("3D Properties:")
                    for prop, value in chart_data['3d_properties'].items():
                        text_lines.append(f"  {prop}: {value}")

            if chart_data['categories']:
                text_lines.append(f"Categories: {len(chart_data['categories'])}")
                text_lines.append("Category: Value (Percentage)")
                if chart_data['series_data']:
                    for series in chart_data['series_data']:
                        series_name = series.get('name', 'Unnamed Series')
                        values = series.get('values', [])
                        if values and chart_data['categories']:
                            total = sum(values) if sum(values) > 0 else 1
                            for cat, val in zip(chart_data['categories'], values):
                                pct = round((val / total) * 100, 2) if val else 0
                                formatted_val = f"{int(val):,}" if val >= 1000 else f"{val}"
                                text_lines.append(f"{cat}: {formatted_val} ({pct}%)")

            text_lines.append("CHART END")
            chart_data['text_representation'] = text_lines

            # Detect table-bar hybrid and extract labels (from original code, assumes first column in Excel)
            if 'bar' in chart_data['chart_type'].lower() and chart_data.get('excel_data'):
                labels = []
                excel = chart_data['excel_data']
                if excel and 'data' in excel and len(excel['data']) > 1:
                    for row in excel['data'][1:]:
                        if row and len(row) > 0:
                            labels.append(row[0])
                chart_data['table_bar_labels'] = labels

            return chart_data

        except Exception as e:
            print(f"    ⚠️  Error processing chart: {e}")
            return None

    def _extract_chart_data_from_excel(self, chart_data: dict, excel_data: dict):
        """
        Extracts chart-specific data (categories, series) from a matched Excel sheet.
        Adds Z-values for 3D charts.
        """
        try:
            rows = excel_data.get('data', [])
            analysis = excel_data.get('analysis', {})

            if not rows or len(rows) < 2:
                return

            headers = rows[0]
            data_rows = rows[1:]

            if analysis.get('potential_categories'):
                chart_data['categories'] = analysis['potential_categories'][:20]
            elif data_rows:
                chart_data['categories'] = [row[0] if row else f"Item {i+1}"
                                          for i, row in enumerate(data_rows)]

            numeric_columns = analysis.get('numeric_columns', [])

            for col_idx in numeric_columns:
                if col_idx < len(headers):
                    series_name = headers[col_idx] if headers[col_idx] else f"Series {col_idx + 1}"
                    values = []
                    for row in data_rows:
                        if col_idx < len(row):
                            try:
                                val = float(str(row[col_idx]).replace(',', ''))
                            except ValueError:
                                val = 0
                            values.append(val)
                        else:
                            values.append(0)

                    series_data = {
                        'name': series_name,
                        'values': values,
                        'column_index': col_idx
                    }

                    if chart_data['is_3d']:
                        series_data['z_values'] = self._generate_3d_values(values, chart_data['chart_type'])

                    chart_data['series_data'].append(series_data)
        except Exception as e:
            chart_data['excel_extraction_error'] = str(e)

    def _generate_3d_values(self, values: list, chart_type: str) -> list:
        """
        Generates appropriate Z-values for 3D charts based on their type.
        """
        if not values:
            return []

        chart_type_lower = chart_type.lower()
        strategies = {
            'column3d': lambda v: [1.0] * len(v),
            'bar3d': lambda v: [1.0] * len(v),
            'line3d': lambda v: [float(i) * 0.5 for i in range(len(v))],
            'pie3d': lambda v: [0.5] * len(v),
            'area3d': lambda v: [float(val) * 0.1 for val in v],
            'surface3d': lambda v: [float(val) * 0.2 for val in v],
        }

        for key, strategy in strategies.items():
            if key in chart_type_lower:
                return strategy(values)

        return [float(i) + (val * 0.01) for i, val in enumerate(values)]

# --- 6. TableDataProcessor Class ---
class TableDataProcessor:
    """
    Handles the extraction of data from PowerPoint table objects.
    """
    def __init__(self, output_dirs: dict, user_login: str, current_time: datetime):
        """
        Initializes the table data processor.

        Args:
            output_dirs (dict): Dictionary of output directories for saving table data.
            user_login (str): User identifier for metadata.
            current_time (datetime): Timestamp for metadata.
        """
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time

    def process_table_data(self, table_shape, slide_num: int, table_counter: int) -> dict:
        """
        Processes a single PowerPoint table shape, extracting its content
        and saving it as JSON, CSV, and XLSX.

        Args:
            table_shape: The PowerPoint table shape object.
            slide_num (int): The slide number where the table is located.
            table_counter (int): The current table count (for unique IDs).

        Returns:
            dict: A dictionary containing all extracted table data, or None if error.
        """
        try:
            table_id = f"table_{slide_num}_{table_counter}"
            print(f"    📊 Table found: {table_id}")

            table = table_shape.table
            table_data = {
                'table_id': table_id,
                'slide_number': slide_num,
                'shape_index': table_counter, # This was originally shape_idx, mapping to current counter
                'rows': len(table.rows),
                'columns': len(table.columns) if len(table.rows) > 0 else 0,
                'headers': [],
                'data': [],
                'text_representation': [],
                'extraction_metadata': {
                    'extracted_by': self.user_login,
                    'extraction_time': self.current_time.isoformat()
                }
            }

            text_lines = ["TABLE START"]

            if len(table.rows) > 0:
                headers = []
                for cell in table.rows[0].cells:
                    header_text = cell.text.strip()
                    headers.append(header_text)

                table_data['headers'] = headers
                text_lines.append("\t".join(headers))

                for row in table.rows[1:]:
                    row_data = {}
                    row_text = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        row_text.append(cell_text)
                        key = headers[col_idx] if col_idx < len(headers) else f"Column {col_idx + 1}"
                        row_data[key] = cell_text
                    table_data['data'].append(row_data)
                    text_lines.append("\t".join(row_text))

            text_lines.append("TABLE END")
            table_data['text_representation'] = text_lines

            table_json_path = os.path.join(self.output_dirs['tables'], f"{table_id}.json")
            with open(table_json_path, 'w', encoding='utf-8') as f:
                json.dump(table_data, f, indent=2, ensure_ascii=False)

            if table_data['data']:
                df = pd.DataFrame(table_data['data'])
                table_csv_path = os.path.join(self.output_dirs['tables'], f"{table_id}.csv")
                df.to_csv(table_csv_path, index=False, encoding='utf-8')
                table_xlsx_path = os.path.join(self.output_dirs['tables'], f"{table_id}.xlsx")
                df.to_excel(table_xlsx_path, index=False)

            return table_data

        except Exception as e:
            print(f"    ⚠️  Error processing table: {e}")
            return None

# --- 7. SlideContentProcessor Class ---
class SlideContentProcessor:
    """
    Responsible for iterating through slides, extracting text, tables, and charts,
    and then coordinating with the ChartDataProcessor and TableDataProcessor.
    """
    def __init__(self, pptx_path: str, output_dirs: dict, user_login: str, current_time: datetime,
                 chart_data_processor: ChartDataProcessor, table_data_processor: TableDataProcessor,
                 chart_visualizer: ChartVisualizer, doc_id: str):
        """
        Initializes the slide content processor.

        Args:
            pptx_path (str): Path to the PowerPoint file.
            output_dirs (dict): Dictionary of output directories.
            user_login (str): User identifier for metadata.
            current_time (datetime): Timestamp for metadata.
            chart_data_processor (ChartDataProcessor): Instance for chart data processing.
            table_data_processor (TableDataProcessor): Instance for table data processing.
            chart_visualizer (ChartVisualizer): Instance for chart visualization.
            doc_id (str): A unique ID for the document (e.g., base filename) for image naming.
        """
        self.pptx_path = pptx_path
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time
        self.chart_data_processor = chart_data_processor
        self.table_data_processor = table_data_processor
        self.chart_visualizer = chart_visualizer
        self.doc_id = doc_id

        self.all_charts = []
        self.all_3d_charts = []
        self.chart_errors = []
        self.table_errors = []
        self.chart_counter = 0
        self.table_counter = 0
        self.chart_3d_counter = 0

    def process_presentation(self):
        """
        Loads the PowerPoint presentation and iterates through each slide
        to extract text, tables, and charts.
        """
        print(f"\n📖 Processing PowerPoint presentation...")
        prs = Presentation(self.pptx_path)
        print(f"📄 Found {len(prs.slides)} slides to process.")

        for slide_idx, slide in enumerate(prs.slides, start=1):
            print(f"\n🔍 Processing Slide {slide_idx}...")

            slide_data = {
                'slide_number': slide_idx,
                'slide_title': self._extract_slide_title(slide, slide_idx),
                'text_content': [],
                'charts': [],
                'tables': [],
                'extraction_metadata': {
                    'extracted_by': self.user_login,
                    'extraction_time': self.current_time.isoformat()
                }
            }

            slide_text_lines = [f"[SLIDE {slide_idx}: {slide_data['slide_title']}]"]
            slide_text_lines.append(f"Extracted by: {self.user_login}")
            slide_text_lines.append(f"Extraction time: {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC")
            slide_text_lines.append("=" * 60)

            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        text = shape.text_frame.text.strip()
                        slide_data['text_content'].append(text)
                        slide_text_lines.append(f"\n[TEXT {shape_idx + 1}]")
                        slide_text_lines.append(text)

                    if shape.has_table:
                        self.table_counter += 1
                        table_data = self.table_data_processor.process_table_data(shape, slide_idx, self.table_counter)
                        if table_data:
                            slide_data['tables'].append(table_data)
                            slide_text_lines.append(f"\n[TABLE {len(slide_data['tables'])}]")
                            slide_text_lines.extend(table_data['text_representation'])
                        else:
                            self.table_errors.append({'slide': slide_idx, 'shape': shape_idx, 'error': 'Table processing failed'})


                    if hasattr(shape, 'chart'):
                        self.chart_counter += 1
                        chart_data = self.chart_data_processor.process_chart_data(shape, slide_idx, self.chart_counter -1 ) #excel_idx used chart_counter-1 for 0-indexing
                        if chart_data:
                            slide_data['charts'].append(chart_data)

                            if chart_data.get('is_3d'):
                                self.chart_3d_counter += 1
                                self.all_3d_charts.append(chart_data)
                                slide_text_lines.append(f"\n[3D CHART {self.chart_3d_counter}]")
                            else:
                                slide_text_lines.append(f"\n[CHART {self.chart_counter}]")

                            slide_text_lines.extend(chart_data['text_representation'])
                            self.all_charts.append(chart_data)

                            # Create chart visualization and potentially save image
                            self.chart_visualizer.create_chart_visualization(
                                chart_data, self.chart_counter, self.doc_id, slide_idx
                            )
                            # Attempt to save the raw chart image if available
                            try:
                                if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.image:
                                    chart_img_path = os.path.join(
                                        self.output_dirs['charts_3d'] if chart_data['is_3d'] else self.output_dirs['charts_2d'],
                                        f"{self.doc_id}_slide{slide_idx}_chart{self.chart_counter}.png"
                                    )
                                    with open(chart_img_path, 'wb') as f:
                                        f.write(shape.image.blob)
                                    chart_data['image_path'] = chart_img_path
                                    print(f"      🖼️ Saved raw chart image: {os.path.basename(chart_img_path)}")
                            except Exception as img_e:
                                print(f"      ⚠️  Could not save raw chart image for {chart_id}: {img_e}")
                                chart_data['image_path'] = None # Ensure it's explicitly None on failure
                        else:
                            self.chart_errors.append({'slide': slide_idx, 'shape': shape_idx, 'error': 'Chart processing failed'})

                except Exception as e:
                    error_info = {
                        'slide': slide_idx,
                        'shape': shape_idx,
                        'error': str(e),
                        'error_type': 'shape_processing'
                    }
                    self.chart_errors.append(error_info)
                    print(f"    ⚠️  Error processing shape {shape_idx}: {e}")

            slide_json_path = os.path.join(self.output_dirs['slide_text'], f"slide_{slide_idx}_data.json")
            with open(slide_json_path, 'w', encoding='utf-8') as f:
                json.dump(slide_data, f, indent=2, ensure_ascii=False)

            slide_text_path = os.path.join(self.output_dirs['slide_text'], f"slide_{slide_idx}.txt")
            with open(slide_text_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(slide_text_lines))

            print(f"  📄 Slide {slide_idx}: {len(slide_data['charts'])} charts, {len(slide_data['tables'])} tables.")

    def _extract_slide_title(self, slide, slide_num: int) -> str:
        """
        Extracts the slide title, with fallbacks if no explicit title shape is found.
        """
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                if len(text) < 100 and '\n' not in text:
                    return text
        return f"Slide {slide_num}"

    def get_all_charts(self) -> list:
        return self.all_charts

    def get_all_3d_charts(self) -> list:
        return self.all_3d_charts

    def get_chart_count(self) -> int:
        return self.chart_counter

    def get_3d_chart_count(self) -> int:
        return self.chart_3d_counter

    def get_table_count(self) -> int:
        return self.table_counter

    def get_chart_errors(self) -> list:
        return self.chart_errors

    def get_table_errors(self) -> list:
        return self.table_errors

# --- 8. ReportGenerator Class ---
class ReportGenerator:
    """
    Responsible for generating summary and detailed analysis reports
    in JSON and Markdown formats.
    """
    def __init__(self, output_dirs: dict, user_login: str, current_time: datetime, pptx_filename: str):
        """
        Initializes the report generator.

        Args:
            output_dirs (dict): Dictionary of output directories for saving reports.
            user_login (str): User identifier for reports.
            current_time (datetime): Timestamp for reports.
            pptx_filename (str): Original PowerPoint filename for reports.
        """
        self.output_dirs = output_dirs
        self.user_login = user_login
        self.current_time = current_time
        self.pptx_filename = pptx_filename

    def generate_analysis_reports(self, chart_stats: dict, excel_data: list, xml_charts_analysis: dict, errors: dict):
        """
        Generates comprehensive analysis reports based on extracted data.

        Args:
            chart_stats (dict): Statistics and lists of charts (total, 3D, 2D).
            excel_data (list): List of extracted embedded Excel data.
            xml_charts_analysis (dict): Dictionary of analyzed XML chart data.
            errors (dict): Dictionary containing chart_errors and table_errors lists.
        """
        print(f"\n📊 Generating analysis reports...")
        self._generate_main_summary(chart_stats, xml_charts_analysis, excel_data, errors)
        if chart_stats['all_3d_charts']:
            self._generate_3d_charts_analysis(chart_stats['all_3d_charts'])
        if excel_data:
            self._generate_excel_analysis(excel_data)
        if errors['chart_errors'] or errors['table_errors']:
            self._generate_error_reports(errors)

    def _generate_main_summary(self, chart_stats: dict, xml_charts_analysis: dict, excel_data: list, errors: dict):
        """Generates the main summary report in JSON and Markdown."""
        summary_data = {
            'extraction_summary': {
                'source_file': self.pptx_filename,
                'extracted_by': self.user_login,
                'extraction_time': self.current_time.isoformat(),
                'extractor_version': 'enhanced_3d_v2.0_refactored'
            },
            'statistics': {
                'total_charts': chart_stats['chart_counter'],
                'charts_3d': chart_stats['chart_3d_counter'],
                'charts_2d': chart_stats['chart_counter'] - chart_stats['chart_3d_counter'],
                'tables': chart_stats['table_counter'],
                'excel_files': len(excel_data),
                'xml_charts_analyzed': len(xml_charts_analysis),
                'errors': len(errors['chart_errors']) + len(errors['table_errors'])
            },
            'charts_breakdown': {
                'chart_types': {},
                '3d_chart_types': {},
                'charts_with_excel_data': 0
            }
        }

        for chart in chart_stats['all_charts']:
            chart_type = chart.get('chart_type', 'Unknown')
            summary_data['charts_breakdown']['chart_types'][chart_type] = \
                summary_data['charts_breakdown']['chart_types'].get(chart_type, 0) + 1
            if chart.get('is_3d'):
                summary_data['charts_breakdown']['3d_chart_types'][chart_type] = \
                    summary_data['charts_breakdown']['3d_chart_types'].get(chart_type, 0) + 1
            if chart.get('excel_data'):
                summary_data['charts_breakdown']['charts_with_excel_data'] += 1

        summary_path_json = os.path.join(self.output_dirs['analysis'], 'extraction_summary.json')
        with open(summary_path_json, 'w', encoding='utf-8') as f:
            json.dump(summary_data, f, indent=2, ensure_ascii=False)

        md_path = os.path.join(self.output_dirs['analysis'], 'extraction_summary.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(f"# Enhanced PowerPoint Extraction Summary\n\n")
            f.write(f"**Extracted by:** {self.user_login}\n")
            f.write(f"**Extraction Time:** {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC\n")
            f.write(f"**Source File:** {self.pptx_filename}\n\n")

            f.write(f"## 📊 Statistics\n\n")
            f.write(f"- **Total Charts:** {chart_stats['chart_counter']}\n")
            f.write(f"- **3D Charts:** {chart_stats['chart_3d_counter']}\n")
            f.write(f"- **2D Charts:** {chart_stats['chart_counter'] - chart_stats['chart_3d_counter']}\n")
            f.write(f"- **Tables:** {chart_stats['table_counter']}\n")
            f.write(f"- **Excel Files:** {len(excel_data)}\n")
            f.write(f"- **XML Charts Analyzed:** {len(xml_charts_analysis)}\n")
            f.write(f"- **Errors Encountered:** {len(errors['chart_errors']) + len(errors['table_errors'])}\n\n")

            if chart_stats['chart_3d_counter'] > 0:
                f.write(f"## 🎲 3D Charts Found\n\n")
                for i, chart in enumerate(chart_stats['all_3d_charts'], 1):
                    f.write(f"{i}. **{chart['chart_id']}** (Slide {chart['slide_number']})\n")
                    f.write(f"   - Type: {chart['chart_type']}\n")
                    if chart.get('3d_properties'):
                        f.write(f"   - 3D Properties: {len(chart['3d_properties'])} found\n")
                    f.write(f"\n")
        print(f"  📄 Main summary saved: {md_path}")

    def _generate_3d_charts_analysis(self, all_3d_charts: list):
        """Generates detailed 3D charts analysis."""
        analysis_data = {
            'total_3d_charts': len(all_3d_charts),
            'analysis_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            },
            'charts': all_3d_charts,
            '3d_properties_summary': {},
            'chart_types_distribution': {}
        }

        all_properties = {}
        for chart in all_3d_charts:
            chart_type = chart.get('chart_type', 'Unknown')
            analysis_data['chart_types_distribution'][chart_type] = \
                analysis_data['chart_types_distribution'].get(chart_type, 0) + 1

            for prop, value in chart.get('3d_properties', {}).items():
                if prop not in all_properties:
                    all_properties[prop] = []
                all_properties[prop].append(value)

        analysis_data['3d_properties_summary'] = {
            prop: {'count': len(values), 'unique_values': list(set(values))}
            for prop, values in all_properties.items()
        }

        analysis_path = os.path.join(self.output_dirs['analysis'], '3d_charts_analysis.json')
        with open(analysis_path, 'w', encoding='utf-8') as f:
            json.dump(analysis_data, f, indent=2, ensure_ascii=False)
        print(f"  🎲 3D charts analysis saved: {analysis_path}")

    def _generate_excel_analysis(self, embedded_data: list):
        """Generates Excel data analysis."""
        excel_analysis = {
            'total_files': len(embedded_data),
            'analysis_metadata': {
                'analyzed_by': self.user_login,
                'analysis_time': self.current_time.isoformat()
            },
            'files': embedded_data,
            'data_structure_summary': {},
            'sheets_summary': {
                'total_sheets': 0,
                'chart_suitable_sheets': 0,
                'multi_series_sheets': 0
            }
        }

        for excel_file in embedded_data:
            for sheet in excel_file.get('sheets', []):
                excel_analysis['sheets_summary']['total_sheets'] += 1
                structure = sheet.get('analysis', {}).get('data_structure', 'unknown')
                excel_analysis['data_structure_summary'][structure] = \
                    excel_analysis['data_structure_summary'].get(structure, 0) + 1
                if structure == 'chart_suitable':
                    excel_analysis['sheets_summary']['chart_suitable_sheets'] += 1
                elif structure == 'multi_series':
                    excel_analysis['sheets_summary']['multi_series_sheets'] += 1

        analysis_path = os.path.join(self.output_dirs['analysis'], 'excel_data_analysis.json')
        with open(analysis_path, 'w', encoding='utf-8') as f:
            json.dump(excel_analysis, f, indent=2, ensure_ascii=False)
        print(f"  📊 Excel analysis saved: {analysis_path}")

    def _generate_error_reports(self, errors: dict):
        """Generates error reports in Excel format."""
        if errors['chart_errors']:
            errors_df = pd.DataFrame(errors['chart_errors'])
            errors_path = os.path.join(self.output_dirs['errors'], 'chart_errors.xlsx')
            errors_df.to_excel(errors_path, index=False)
            print(f"  ⚠️  Chart errors saved: {errors_path}")

        if errors['table_errors']:
            errors_df = pd.DataFrame(errors['table_errors'])
            errors_path = os.path.join(self.output_dirs['errors'], 'table_errors.xlsx')
            errors_df.to_excel(errors_path, index=False)
            print(f"  ⚠️  Table errors saved: {errors_path}")

# --- 9. MarkdownExporter Class ---
class MarkdownExporter:
    """
    Dedicated to saving extracted slide content as Markdown files (one file per slide).
    """
    def __init__(self, output_dirs: dict):
        """
        Initializes the Markdown exporter.

        Args:
            output_dirs (dict): Dictionary of output directories, specifically for 'slide_text'
                                and 'markdown_slides'.
        """
        self.output_dirs = output_dirs

    def save_slides_as_markdown(self):
        """
        Reads the processed slide data (JSON) and converts it into
        structured Markdown files, one per slide.
        """
        print("\n📝 Saving slides as Markdown...")
        md_dir = self.output_dirs['markdown_slides']
        os.makedirs(md_dir, exist_ok=True)

        slide_files = sorted([
            f for f in os.listdir(self.output_dirs['slide_text'])
            if f.startswith("slide_") and f.endswith("_data.json")
        ], key=lambda x: int(x.split('_')[1]))

        for slide_json in slide_files:
            slide_num = slide_json.split('_')[1]
            with open(os.path.join(self.output_dirs['slide_text'], slide_json), encoding='utf-8') as f:
                slide_data = json.load(f)

            lines = [f"# Slide {slide_num}: {slide_data.get('slide_title', '')}\n"]

            for text in slide_data.get('text_content', []):
                lines.append(text)
                lines.append("")

            for table in slide_data.get('tables', []):
                lines.append("## Table")
                if 'headers' in table and table['headers']:
                    lines.append('| ' + ' | '.join(table['headers']) + ' |')
                    lines.append('| ' + ' | '.join(['---'] * len(table['headers'])) + ' |')
                    for row in table.get('data', []):
                        row_vals = [str(row.get(h, '')) for h in table['headers']]
                        lines.append('| ' + ' | '.join(row_vals) + ' |')
                lines.append("")

            for chart in slide_data.get('charts', []):
                lines.append(f"## Chart ({chart.get('chart_type', '')})")
                if chart.get('is_3d'):
                    lines.append("### 🎲 3D Chart")
                if chart.get('categories') and chart.get('series_data'):
                    for series in chart['series_data']:
                        lines.append(f"**{series.get('name', 'Series')}**")
                        if chart['categories']:
                            lines.append('| Category | Value |')
                            lines.append('|---|---|')
                            for cat, val in zip(chart['categories'], series.get('values', [])):
                                lines.append(f'| {cat} | {val} |')
                        lines.append("")
                if chart.get('visualization_path'):
                    lines.append(f"![Chart Visualization]({os.path.basename(chart.get('visualization_path'))})") # Referencing filename only
                    lines.append("")
                lines.append("")

            md_path = os.path.join(md_dir, f"slide_{slide_num}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(lines))
            print(f"✅ Saved Markdown: {md_path}")

# --- 10. PowerPointExtractor (Main Orchestrator) Class ---
class PowerPointExtractor:
    """
    Main orchestrator class for extracting data from PowerPoint presentations.
    It coordinates the activities of specialized helper classes.
    """
    def __init__(self, pptx_path: str = None, user_login: str = "Rateb1223"):
        """
        Initializes the main PowerPoint Extractor.

        Args:
            pptx_path (str): Path to the PowerPoint file. If None, it auto-detects.
            user_login (str): User identifier for tracking and output.
        """
        self.user_login = user_login
        self.current_time = datetime(2025, 6, 3, 16, 16, 32) # Fixed for reproducibility from original code

        print(f"\n{'='*60}")
        print(f"🚀 Initializing Enhanced PowerPoint Extractor")
        print(f"👤 User: {self.user_login}")
        print(f"📅 Current Time: {self.current_time.strftime('%Y-%m-%d %H:%M:%S')} UTC")
        print(f"{'='*60}")

        # Initialize core components
        self.file_handler = PowerPointFileHandler(pptx_path, self.user_login, self.current_time)
        self.output_dirs = self.file_handler.get_paths()
        self.pptx_path = self.file_handler.get_pptx_path()
        self.unzipped_dir = self.file_handler.get_unzipped_dir()
        self.doc_id = Path(self.pptx_path).stem # Used for naming output files

        self.excel_data_extractor = ExcelDataExtractor(self.unzipped_dir, self.output_dirs, self.user_login, self.current_time)
        self.chart_xml_analyzer = ChartXmlAnalyzer(self.unzipped_dir, self.output_dirs, self.user_login, self.current_time)
        self.chart_visualizer = ChartVisualizer(self.output_dirs, self.user_login, self.current_time)

        # These will be initialized after initial data extraction steps
        self.chart_data_processor = None
        self.table_data_processor = TableDataProcessor(self.output_dirs, self.user_login, self.current_time)
        self.slide_content_processor = None
        self.report_generator = ReportGenerator(self.output_dirs, self.user_login, self.current_time, os.path.basename(self.pptx_path))
        self.markdown_exporter = MarkdownExporter(self.output_dirs)

    def extract_all(self) -> str:
        """
        Main extraction method that orchestrates the entire process:
        1. Unzips the PowerPoint.
        2. Extracts embedded Excel data.
        3. Extracts and analyzes chart XML for 3D properties.
        4. Processes slide content (text, tables, charts).
        5. Generates analysis reports.
        6. Saves slides as Markdown.
        7. Cleans up temporary files.

        Returns:
            str: The path to the main output directory, or None if extraction fails.
        """
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            print(f"❌ PowerPoint file not found: {self.pptx_path}")
            return None

        print(f"\n{'='*60}")
        print(f"🔄 Starting Enhanced Extraction Process")
        print(f"{'='*60}")

        try:
            self.file_handler.unzip_pptx()

            self.excel_data_extractor.extract_embedded_excel_data()
            embedded_excel_data = self.excel_data_extractor.get_extracted_data()

            self.chart_xml_analyzer.extract_xml_chart_data()
            xml_charts_analysis = self.chart_xml_analyzer.get_xml_charts_analysis()

            # Initialize processors that depend on prior extraction results
            self.chart_data_processor = ChartDataProcessor(
                embedded_excel_data, xml_charts_analysis, self.user_login, self.current_time
            )
            self.slide_content_processor = SlideContentProcessor(
                self.pptx_path, self.output_dirs, self.user_login, self.current_time,
                self.chart_data_processor, self.table_data_processor, self.chart_visualizer, self.doc_id
            )

            self.slide_content_processor.process_presentation()

            # Collect all stats and errors for reporting
            chart_stats = {
                'all_charts': self.slide_content_processor.get_all_charts(),
                'all_3d_charts': self.slide_content_processor.get_all_3d_charts(),
                'chart_counter': self.slide_content_processor.get_chart_count(),
                'chart_3d_counter': self.slide_content_processor.get_3d_chart_count(),
                'table_counter': self.slide_content_processor.get_table_count(),
            }
            errors = {
                'chart_errors': self.slide_content_processor.get_chart_errors() + self.excel_data_extractor.get_errors(),
                'table_errors': self.slide_content_processor.get_table_errors(),
            }

            self.report_generator.generate_analysis_reports(
                chart_stats, embedded_excel_data, xml_charts_analysis, errors
            )

            self.markdown_exporter.save_slides_as_markdown()

            # The original code's `save_slides_as_markdown` also handled moving files
            # to a global output folder in __main__. We will move that logic to the __main__ block
            # or a separate utility if needed outside the extractor itself.

            print(f"\n✅ Enhanced extraction complete!")
            print(f"📊 Total charts: {chart_stats['chart_counter']}")
            print(f"🎲 3D charts: {chart_stats['chart_3d_counter']}")
            print(f"📋 Tables: {chart_stats['table_counter']}")
            print(f"📁 Excel files: {self.excel_data_extractor.get_excel_count()}")
            print(f"📂 Output: {self.output_dirs['main']}")

            return self.output_dirs['main']

        except Exception as e:
            print(f"❌ Error during extraction: {e}")
            import traceback
            traceback.print_exc()
            return None
        finally:
            self.file_handler.cleanup_temp_files()


# Main execution function for external calls
def run_powerpoint_extraction(pptx_path: str = None, user_login: str = "Rateb1223") -> str:
    """
    Runs the comprehensive PowerPoint data extraction process.

    Args:
        pptx_path (str): Path to the PowerPoint file. If None, it attempts to auto-detect.
        user_login (str): User identifier for tracking and output.

    Returns:
        str: The path to the output directory if successful, None otherwise.
    """
    extractor = PowerPointExtractor(pptx_path, user_login)
    return extractor.extract_all()

# Enhanced demo execution
if __name__ == "__main__":
    # Current context (these would typically come from environment or config)
    # Using fixed values for reproducibility as in the original __main__ block
    current_time_for_demo = datetime(2025, 6, 3, 16, 16, 32)
    user_login_for_demo = "Rateb1223"

    print(f"🚀 Enhanced 3D PowerPoint Extractor Demo (Refactored)")
    print(f"👤 User: {user_login_for_demo}")
    print(f"📅 Current Time: {current_time_for_demo.strftime('%Y-%m-%d %H:%M:%S')} UTC")
    print(f"📍 Working Directory: {os.getcwd()}")

    # Assuming 'pptx_input_dir' exists and contains PowerPoint files for the demo
    pptx_input_dir = 'OlmOCR/Input_PDFs'
    if not os.path.exists(pptx_input_dir):
        print(f"Warning: Demo input directory '{pptx_input_dir}' not found. Please create it and add .pptx files.")
        pptx_input_dir = '.' # Fallback to current directory for auto-detection

    pptx_files = [os.path.join(pptx_input_dir, f)
                  for f in os.listdir(pptx_input_dir)
                  if f.endswith(('.pptx', '.ppt', '.potx'))]

    # Global markdown output folder for all processed PPTX files
    markdown_out_dir = os.path.join(os.getcwd(), "pptx_markdown_split_slides")
    os.makedirs(markdown_out_dir, exist_ok=True)

    if not pptx_files:
        print("\n❌ No PowerPoint files found in input directory or current directory.")
        print("📝 Please add PowerPoint files to test the enhanced extractor.")
        print("\n📋 Features of this enhanced extractor:")
        print("  ✅ Advanced 3D chart detection and analysis")
        print("  ✅ XML-based 3D properties extraction")
        print("  ✅ Comprehensive Excel data matching")
        print("  ✅ Enhanced visualizations with 3D support")
        print("  ✅ Structured output with analysis reports")
        print("  ✅ Error handling and logging")
        print("  ✅ Dynamic paths and user tracking")
    else:
        print(f"\n🔍 Found {len(pptx_files)} PowerPoint files:")
        for i, file in enumerate(pptx_files, 1):
            print(f"  {i}. {file}")

        for pptx_file in pptx_files:
            print(f"\n📄 Processing file: {pptx_file}")
            # Use the refactored main function
            output_folder = run_powerpoint_extraction(pptx_file, user_login_for_demo)

            if output_folder:
                # Move/copy markdown files to the global output folder
                src_md_dir = os.path.join(output_folder, "markdown_slides")
                dst_md_dir = os.path.join(markdown_out_dir, Path(pptx_file).stem)
                os.makedirs(dst_md_dir, exist_ok=True)
                for md_file in os.listdir(src_md_dir):
                    src_file = os.path.join(src_md_dir, md_file)
                    dst_file = os.path.join(dst_md_dir, md_file)
                    shutil.copy2(src_file, dst_file)
                print(f"✅ Markdown slides for {Path(pptx_file).name} copied to {dst_md_dir}")
            else:
                print(f"❌ Extraction failed for {Path(pptx_file).name}.")

        print(f"\n🎉 All slides from processed files exported as markdown in: {markdown_out_dir}")
